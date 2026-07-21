from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # 진한 네이비
C_ACCENT    = RGBColor(0x00, 0xB4, 0xD8)   # 하늘색 포인트
C_ACCENT2   = RGBColor(0x90, 0xE0, 0xEF)   # 연한 하늘
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCA, 0xF0, 0xF8)
C_DARK_BOX  = RGBColor(0x1B, 0x2E, 0x42)   # 카드 배경
C_GRAY      = RGBColor(0xAA, 0xBB, 0xCC)
C_GREEN     = RGBColor(0x06, 0xD6, 0xA0)
C_ORANGE    = RGBColor(0xFF, 0x9F, 0x1C)
C_RED       = RGBColor(0xEF, 0x47, 0x6F)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]   # 완전 빈 레이아웃

# ══════════════════════════════════════════════════════════════
# 헬퍼 함수들
# ══════════════════════════════════════════════════════════════
def add_bg(slide, color=C_BG):
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width: shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=14, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, l, t, w, h, size=12, color=C_WHITE,
                  bold_first=False, line_spacing=None):
    """lines: list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, bld, col = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            col = item[2] if len(item) > 2 else color
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = bld
        run.font.color.rgb = col
    return txb

def accent_bar(slide, t_cm, h_cm=0.08):
    """가로 포인트 라인"""
    bar = slide.shapes.add_shape(1, Cm(1.2), Cm(t_cm),
                                 Cm(30.8), Cm(h_cm))
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

def section_badge(slide, text, l, t, w=3.5, h=0.65,
                  bg=C_ACCENT, fg=C_BG):
    r = add_rect(slide, l, t, w, h, bg)
    add_text(slide, text, l+0.15, t+0.05, w-0.2, h-0.1,
             size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)

def dot_bullet(slide, items, l, t, w, size=11.5, color=C_WHITE,
               dot_color=C_ACCENT, spacing=0.72):
    for i, item in enumerate(items):
        # dot
        d = slide.shapes.add_shape(9, Cm(l), Cm(t + i*spacing + 0.12),
                                   Cm(0.22), Cm(0.22))
        d.fill.solid(); d.fill.fore_color.rgb = dot_color
        d.line.fill.background()
        add_text(slide, item, l+0.38, t + i*spacing, w-0.38,
                 0.65, size=size, color=color)

# ══════════════════════════════════════════════════════════════
# 슬라이드 1 ── 표지
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# 좌측 포인트 바
bar = slide.shapes.add_shape(1, 0, 0, Cm(0.55), prs.slide_height)
bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
bar.line.fill.background()

# 우측 장식 사각형
add_rect(slide, 26.5, 0, 7.0, 19.05, RGBColor(0x11,0x26,0x3A))
add_rect(slide, 28.0, 3.5, 4.0, 4.0, C_ACCENT)

# 상단 작은 태그
section_badge(slide, "2026 연구사업 계획", 1.5, 1.0, 4.5, 0.55,
              bg=C_ACCENT, fg=C_BG)

# 메인 타이틀
add_text(slide, "도시교통 디지털트윈", 1.5, 2.0, 22, 1.8,
         size=44, bold=True, color=C_WHITE)
add_text(slide, "시뮬레이션 플랫폼 고도화", 1.5, 3.5, 22, 1.6,
         size=40, bold=True, color=C_ACCENT)

add_text(slide, "네트워크 편집 · 실시간 시뮬레이션 · 지능형 분석 · 운영 자동화",
         1.5, 5.1, 23, 0.7, size=15, color=C_ACCENT2)

accent_bar(slide, 5.95, 0.06)

add_text(slide, "한국도로공사 도로교통연구원   |   IITP 지원과제",
         1.5, 6.2, 18, 0.55, size=11, color=C_GRAY)
add_text(slide, "2026. 03", 1.5, 6.75, 6, 0.5, size=11, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 2 ── 목차
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)

# 헤더 밴드
add_rect(slide, 0, 0, 33.87, 2.2, C_DARK_BOX)
accent_bar(slide, 2.2, 0.07)
add_text(slide, "Contents", 1.2, 0.3, 20, 1.0,
         size=30, bold=True, color=C_ACCENT)
add_text(slide, "연구 추진 체계 및 세부 과제", 1.2, 1.25, 20, 0.7,
         size=13, color=C_GRAY)

chapters = [
    ("01", "연구 배경 및 필요성",       "현황과 한계, 추진 근거"),
    ("02", "기존 플랫폼 현황",           "핵심 구현 기능 요약"),
    ("03", "네트워크 편집 고도화",       "위상 검증 · OSM 임포트 · Diff 시각화"),
    ("04", "시뮬레이션 고도화",         "실시간 연동 · 이벤트 편집 · A/B 비교"),
    ("05", "분석 기능 확장",             "교차로 성능 · 속도 프로파일 · 탄소 배출"),
    ("06", "UX 및 운영 편의 향상",       "영상 내보내기 · 자동 리포트 · 대시보드"),
    ("07", "추진 일정 및 기대 효과",     "로드맵 · KPI · 활용 방안"),
]

for i, (num, title, desc) in enumerate(chapters):
    row = i % 4; col = i // 4
    lx = 1.2 + col * 16.5
    ty = 2.7 + row * 1.12

    add_rect(slide, lx, ty, 15.5, 0.95, C_DARK_BOX)
    add_rect(slide, lx, ty, 1.1, 0.95, C_ACCENT)
    add_text(slide, num, lx+0.08, ty+0.12, 1.0, 0.7,
             size=15, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    add_text(slide, title, lx+1.25, ty+0.07, 9.5, 0.5,
             size=13, bold=True, color=C_WHITE)
    add_text(slide, desc,  lx+1.25, ty+0.52, 14.0, 0.4,
             size=9.5, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 3 ── 연구 배경 및 필요성
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
add_text(slide, "01  연구 배경 및 필요성", 1.2, 0.35, 25, 1.0,
         size=24, bold=True, color=C_WHITE)

# 3열 카드
cards = [
    (C_ACCENT, "교통 혼잡 비용", [
        "연간 약 73조 원 (2023)",
        "수도권 집중화 지속",
        "C-ITS·자율주행 확산",
        "정밀 교통 분석 수요 급증",
    ]),
    (C_GREEN, "디지털트윈 정책", [
        "국토부 스마트시티 플랫폼",
        "교통 디지털트윈 국가 R&D",
        "공간정보 오픈API 확대",
        "실시간 연계 인프라 구축",
    ]),
    (C_ORANGE, "기존 시스템 한계", [
        "정적 데이터 기반 재생만 가능",
        "편집 후 오류 자동 검출 없음",
        "외부 데이터 임포트 미지원",
        "분석 결과 보고서화 불가",
    ]),
]
for j, (col, title, items) in enumerate(cards):
    lx = 1.2 + j*10.8
    add_rect(slide, lx, 2.4, 10.2, 4.6, C_DARK_BOX)
    add_rect(slide, lx, 2.4, 10.2, 0.6, col)
    add_text(slide, title, lx+0.3, 2.47, 9.5, 0.5,
             size=13, bold=True, color=C_BG)
    dot_bullet(slide, items, lx+0.3, 3.25, 9.5,
               size=11, dot_color=col, spacing=0.75)

add_text(slide, "※ 본 연구는 기존 플랫폼의 핵심 한계를 극복하고, "
                "실시간·지능형 교통 디지털트윈 플랫폼으로 도약하는 것을 목표로 합니다.",
         1.2, 7.1, 31, 0.6, size=10, color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════
# 슬라이드 4 ── 기존 플랫폼 현황
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
add_text(slide, "02  기존 플랫폼 핵심 기능 현황", 1.2, 0.35, 25, 1.0,
         size=24, bold=True, color=C_WHITE)

features = [
    (C_ACCENT,  "도로 네트워크 편집",
     "Node · Link · Lane · Connection\n스키마 기반 동적 UI, 히스토리 관리"),
    (C_GREEN,   "차량 시뮬레이션",
     "CZML + Web Worker 비동기 처리\n2D(OpenLayers) + 3D(Cesium) 동시 뷰"),
    (C_ORANGE,  "교통 제어 편집",
     "신호 타이밍, 노면마킹\n버스·철도 정류장 관리"),
    (C_ACCENT2, "분석 레이어",
     "히트맵, OD 매트릭스\nTrail 분석, Recharts 기반 통계"),
]
for j, (col, title, desc) in enumerate(features):
    lx = 1.2 + j*7.9
    ty = 2.4
    add_rect(slide, lx, ty, 7.4, 4.5, C_DARK_BOX)
    # 상단 컬러 아이콘 바
    add_rect(slide, lx, ty, 7.4, 0.55, col)
    add_text(slide, f"{j+1:02d}", lx+0.2, ty+0.08, 1.0, 0.4,
             size=13, bold=True, color=C_BG)
    add_text(slide, title, lx+1.0, ty+0.08, 6.0, 0.4,
             size=12, bold=True, color=C_BG)
    for di, line in enumerate(desc.split('\n')):
        add_text(slide, "· " + line,
                 lx+0.3, ty+0.85+di*0.75, 6.7, 0.65,
                 size=11, color=C_LIGHT)

# 하단 기술 스택 바
add_rect(slide, 1.2, 7.1, 31.5, 0.55, C_DARK_BOX)
add_text(slide, "Tech Stack  :  Spring Boot  ·  PostgreSQL  ·  React/TypeScript  ·  "
                "CesiumJS  ·  OpenLayers  ·  Web Workers  ·  Recharts",
         1.5, 7.14, 31, 0.45, size=10, color=C_ACCENT2)

# ══════════════════════════════════════════════════════════════
# 슬라이드 5 ── 네트워크 편집 고도화
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
section_badge(slide, "세부과제 1", 1.2, 0.3, 3.0, 0.55)
add_text(slide, "03  네트워크 편집 고도화", 4.5, 0.35, 25, 1.0,
         size=24, bold=True, color=C_WHITE)

tasks_net = [
    (C_ACCENT, "위상 검증 (Topology Validation)", [
        "단절 링크·고립 노드·차선 방향 충돌 자동 감지",
        "중복 좌표·역방향 링크 실시간 검출",
        "오류 요소 지도 하이라이트 + 원클릭 이동",
        "저장 전 검증 게이트 → 데이터 품질 보증",
        "오류 유형별 통계 패널 제공",
    ]),
    (C_GREEN, "OSM / SHP 외부 데이터 임포트", [
        "OpenStreetMap Overpass API 연동",
        "SHP(Shapefile) 파일 업로드 및 변환",
        "좌표계 자동 변환 (WGS84 ↔ EPSG:5186)",
        "임포트 후 자동 위상 정제 파이프라인",
        "현장 조사 없이 초기 네트워크 구축 가능",
    ]),
    (C_ORANGE, "버전 간 네트워크 Diff 시각화", [
        "v1 → v2 변경 사항 색상 구분 (추가/삭제/수정)",
        "링크·노드·차선별 변경 이력 패널",
        "세션 내 이력 → 서버 영구 저장 이력으로 확장",
        "특정 버전 원클릭 복원 기능",
        "변경 요약 리포트 자동 생성",
    ]),
]
for j, (col, title, items) in enumerate(tasks_net):
    lx = 1.2 + j*10.8
    ty = 2.35
    add_rect(slide, lx, ty, 10.3, 4.75, C_DARK_BOX)
    add_rect(slide, lx, ty, 0.35, 4.75, col)
    add_text(slide, title, lx+0.6, ty+0.15, 9.4, 0.6,
             size=12.5, bold=True, color=col)
    dot_bullet(slide, items, lx+0.6, ty+0.85, 9.3,
               size=10.5, dot_color=col, spacing=0.72)

add_text(slide, "기대 효과  :  데이터 품질 40% 향상  ·  초기 구축 시간 60% 단축  ·  변경 추적 체계화",
         1.2, 7.2, 31, 0.5, size=10.5, color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════
# 슬라이드 6 ── 시뮬레이션 고도화
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
section_badge(slide, "세부과제 2", 1.2, 0.3, 3.0, 0.55, bg=C_GREEN, fg=C_BG)
add_text(slide, "04  시뮬레이션 고도화", 4.5, 0.35, 25, 1.0,
         size=24, bold=True, color=C_WHITE)

tasks_sim = [
    (C_ACCENT, "실시간 시뮬레이션 연동", [
        "WebSocket / SSE 기반 스트리밍 수신",
        "SUMO·VISSIM 시뮬레이터 서버 연동",
        "진행 중인 시뮬레이션 라이브 모니터링",
        "재생·정지·배속 실시간 제어",
        "연결 끊김 자동 재연결 및 버퍼 관리",
    ]),
    (C_GREEN, "시나리오 이벤트 편집기", [
        "교통사고·공사구간·신호 고장 타임라인 삽입",
        "드래그 앤 드롭 이벤트 배치 UI",
        "이벤트 발생 전/후 교통 흐름 자동 비교",
        "이벤트 연쇄 효과 시뮬레이션 (도미노 모델)",
        "시나리오 템플릿 저장·공유 기능",
    ]),
    (C_ORANGE, "A/B 시나리오 비교 뷰", [
        "화면 분할(Split-view) 동기 재생",
        "KPI 나란히 표시 (지체·속도·대기행렬)",
        "시나리오 간 차이 히트맵 오버레이",
        "최적 시나리오 자동 추천 알고리즘",
        "비교 결과 PDF/CSV 내보내기",
    ]),
]
for j, (col, title, items) in enumerate(tasks_sim):
    lx = 1.2 + j*10.8
    ty = 2.35
    add_rect(slide, lx, ty, 10.3, 4.75, C_DARK_BOX)
    add_rect(slide, lx, ty, 0.35, 4.75, col)
    add_text(slide, title, lx+0.6, ty+0.15, 9.4, 0.6,
             size=12.5, bold=True, color=col)
    dot_bullet(slide, items, lx+0.6, ty+0.85, 9.3,
               size=10.5, dot_color=col, spacing=0.72)

add_text(slide, "기대 효과  :  라이브 모니터링 실현  ·  정책 시나리오 검증 효율화  ·  의사결정 지원 고도화",
         1.2, 7.2, 31, 0.5, size=10.5, color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════
# 슬라이드 7 ── 분석 기능 확장
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
section_badge(slide, "세부과제 3", 1.2, 0.3, 3.0, 0.55, bg=C_ORANGE, fg=C_BG)
add_text(slide, "05  분석 기능 확장", 4.5, 0.35, 25, 1.0,
         size=24, bold=True, color=C_WHITE)

tasks_ana = [
    (C_ACCENT, "교차로 성능 분석", [
        "포화도(v/c ratio) 자동 산출",
        "평균 지체 및 최대 대기행렬 계산",
        "신호 현시 최적화 제안 엔진",
        "HCM 기준 서비스 수준(LOS) 등급화",
        "교차로 유형별 비교 벤치마크",
    ]),
    (C_GREEN, "구간 통행시간 / 속도 프로파일", [
        "경로 선택 → 시간대별 속도 그래프",
        "Recharts 기반 인터랙티브 시각화",
        "Peak/Off-peak 자동 구분 분석",
        "링크 속도 이상치 감지 알림",
        "FCD 실측 데이터와 시뮬레이션 비교",
    ]),
    (C_ORANGE, "탄소 배출 추정 모듈", [
        "차량 속도 프로파일 → 배출계수 적용",
        "MOVES / COPERT 배출 모델 내장",
        "차종별 CO₂·NOx·PM2.5 산출",
        "정책 시나리오별 저감 효과 비교",
        "교통영향평가 보고서 자동 연동",
    ]),
]
for j, (col, title, items) in enumerate(tasks_ana):
    lx = 1.2 + j*10.8
    ty = 2.35
    add_rect(slide, lx, ty, 10.3, 4.75, C_DARK_BOX)
    add_rect(slide, lx, ty, 0.35, 4.75, col)
    add_text(slide, title, lx+0.6, ty+0.15, 9.4, 0.6,
             size=12.5, bold=True, color=col)
    dot_bullet(slide, items, lx+0.6, ty+0.85, 9.3,
               size=10.5, dot_color=col, spacing=0.72)

add_text(slide, "기대 효과  :  교차로 운영 효율 15% 향상  ·  탄소 저감 정량화  ·  HCM 기준 준수 검증 자동화",
         1.2, 7.2, 31, 0.5, size=10.5, color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════
# 슬라이드 8 ── UX 및 운영 편의
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
section_badge(slide, "세부과제 4", 1.2, 0.3, 3.0, 0.55, bg=C_RED, fg=C_WHITE)
add_text(slide, "06  UX 및 운영 편의 향상", 4.5, 0.35, 25, 1.0,
         size=24, bold=True, color=C_WHITE)

tasks_ux = [
    (C_ACCENT, "시뮬레이션 영상 내보내기", [
        "Cesium scene.postRender 프레임 캡처",
        "ffmpeg.wasm 기반 브라우저 내 MP4 인코딩",
        "구간 선택 (시작~종료 타임코드) 내보내기",
        "해상도·프레임레이트 설정 옵션",
        "발표·보고서·소셜미디어 공유 지원",
    ]),
    (C_GREEN, "대시보드 & 자동 리포트", [
        "시뮬레이션 종료 후 KPI 요약 화면",
        "PDF 자동 생성 (교통영향평가 형식)",
        "차트·지도·수치 통합 레이아웃",
        "정기 배치 리포트 스케줄링",
        "이메일·슬랙 자동 발송 연동",
    ]),
    (C_ORANGE, "협업 및 접근성 강화", [
        "다중 사용자 동시 편집 (CRDT 기반)",
        "역할별 권한 관리 (뷰어/편집자/관리자)",
        "변경 사항 실시간 댓글·리뷰 기능",
        "반응형 UI (태블릿 현장 조사 지원)",
        "키보드 단축키 및 접근성(WCAG 2.1) 준수",
    ]),
]
for j, (col, title, items) in enumerate(tasks_ux):
    lx = 1.2 + j*10.8
    ty = 2.35
    add_rect(slide, lx, ty, 10.3, 4.75, C_DARK_BOX)
    add_rect(slide, lx, ty, 0.35, 4.75, col)
    add_text(slide, title, lx+0.6, ty+0.15, 9.4, 0.6,
             size=12.5, bold=True, color=col)
    dot_bullet(slide, items, lx+0.6, ty+0.85, 9.3,
               size=10.5, dot_color=col, spacing=0.72)

add_text(slide, "기대 효과  :  보고서 작성 시간 70% 단축  ·  현장 활용성 향상  ·  기관 간 협업 체계 구축",
         1.2, 7.2, 31, 0.5, size=10.5, color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════
# 슬라이드 9 ── 추진 일정 (로드맵)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
add_text(slide, "07  연구 추진 일정", 1.2, 0.35, 25, 1.0,
         size=24, bold=True, color=C_WHITE)

# 분기 헤더
quarters = ["Q1 (1~3월)", "Q2 (4~6월)", "Q3 (7~9월)", "Q4 (10~12월)"]
q_colors  = [C_ACCENT, C_GREEN, C_ORANGE, C_RED]
for qi, (q, col) in enumerate(zip(quarters, q_colors)):
    lx = 5.8 + qi*6.8
    add_rect(slide, lx, 2.2, 6.5, 0.6, col)
    add_text(slide, q, lx+0.2, 2.25, 6.1, 0.5,
             size=11.5, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

# 과제 행
roadmap = [
    ("네트워크 편집",  C_ACCENT,
     ["위상 검증 설계·구현", "OSM 임포트 개발", "Diff 시각화", "통합 테스트"]),
    ("시뮬레이션",     C_GREEN,
     ["WebSocket 설계", "실시간 연동·이벤트 편집기", "A/B 비교뷰", "안정화"]),
    ("분석 기능",      C_ORANGE,
     ["교차로 분석 설계", "속도 프로파일", "탄소 배출 모듈", "HCM 검증"]),
    ("UX / 리포트",   C_RED,
     ["영상 내보내기", "대시보드 개발", "PDF 자동화", "배포·운영"]),
]
for ri, (name, col, tasks) in enumerate(roadmap):
    ty = 3.05 + ri * 1.05
    add_rect(slide, 0.5, ty, 5.0, 0.85, C_DARK_BOX)
    add_rect(slide, 0.5, ty, 0.3, 0.85, col)
    add_text(slide, name, 1.0, ty+0.15, 4.3, 0.55,
             size=11, bold=True, color=col)
    for qi, task in enumerate(tasks):
        lx = 5.8 + qi*6.8
        add_rect(slide, lx+0.15, ty+0.08, 6.2, 0.7, col)
        add_text(slide, task, lx+0.35, ty+0.15, 5.8, 0.55,
                 size=10, bold=False, color=C_BG, align=PP_ALIGN.CENTER)

add_text(slide, "▶  중간 점검 : 6월 말   |   최종 성과물 제출 : 12월",
         1.2, 7.2, 31, 0.5, size=10.5, color=C_GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 10 ── 기대 효과 & 결론
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide)
add_rect(slide, 0, 0, 33.87, 2.0, C_DARK_BOX)
accent_bar(slide, 2.0, 0.07)
add_text(slide, "기대 효과 및 활용 방안", 1.2, 0.35, 28, 1.0,
         size=24, bold=True, color=C_WHITE)

kpis = [
    (C_ACCENT,  "기술적 성과",
     ["위상 검증으로 네트워크 데이터 품질 40% 향상",
      "OSM 임포트로 초기 구축 시간 60% 단축",
      "실시간 시뮬레이션 지연 < 500ms 달성",
      "A/B 비교로 정책 의사결정 근거 정량화"]),
    (C_GREEN,   "사회·경제적 효과",
     ["교통 혼잡 비용 절감 정책 수립 지원",
      "탄소 배출 저감 시나리오 정량 제시",
      "보고서 자동화로 행정 비용 70% 절감",
      "스마트시티 플랫폼 연계 기반 마련"]),
    (C_ORANGE,  "활용 방안",
     ["교통영향평가 보고서 자동 생성 시스템",
      "지자체·공공기관 교통운영 의사결정 도구",
      "C-ITS / 자율주행 시뮬레이션 환경 제공",
      "대학·연구기관 교육·연구 플랫폼 개방"]),
]
for j, (col, title, items) in enumerate(kpis):
    lx = 1.2 + j*10.8
    add_rect(slide, lx, 2.35, 10.3, 4.6, C_DARK_BOX)
    add_rect(slide, lx, 2.35, 10.3, 0.6, col)
    add_text(slide, title, lx+0.3, 2.42, 9.7, 0.48,
             size=13, bold=True, color=C_BG)
    dot_bullet(slide, items, lx+0.3, 3.15, 9.7,
               size=11, dot_color=col, spacing=0.77)

# 하단 강조 문구
add_rect(slide, 1.2, 7.05, 31.5, 0.65, C_ACCENT)
add_text(slide, "본 연구를 통해 대한민국 교통 디지털트윈의 표준 플랫폼을 구축합니다",
         1.5, 7.1, 31, 0.55, size=12.5, bold=True,
         color=C_BG, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 저장
# ══════════════════════════════════════════════════════════════
out = "/Users/hskim/Documents/repo/iitp/IITP_2026_연구사업계획.pptx"
prs.save(out)
print(f"저장 완료 → {out}")
