"""
IITP 3차년도 연구계획 PPT — v5  Modern Design
top-bar card style · shadow · square bullets · refined typography
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color Palette ──────────────────────────────────────────────
C_BG      = RGBColor(0xF0, 0xF4, 0xF8)   # cool light blue-grey bg
C_HDR     = RGBColor(0x0D, 0x14, 0x26)   # near-black header
C_SURF    = RGBColor(0xFF, 0xFF, 0xFF)
C_SURF2   = RGBColor(0xF8, 0xFA, 0xFF)
C_SHADOW  = RGBColor(0xD8, 0xE3, 0xEE)   # card shadow
C_BORDER  = RGBColor(0xCB, 0xD5, 0xE1)
C_DIV     = RGBColor(0xE2, 0xE8, 0xF0)
C_TEXT    = RGBColor(0x0D, 0x14, 0x26)
C_TEXT2   = RGBColor(0x47, 0x55, 0x69)
C_TEXT3   = RGBColor(0x94, 0xA3, 0xB8)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
# Blues
C_BLUE    = RGBColor(0x1D, 0x4E, 0xD8)   # deeper blue
C_BLUE_M  = RGBColor(0x3B, 0x82, 0xF6)
C_BLUE_L  = RGBColor(0xDB, 0xE9, 0xFE)
C_BLUE_D  = RGBColor(0x1E, 0x3A, 0x8A)
# Others
C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
C_AMBER_L = RGBColor(0xFE, 0xF3, 0xC7)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_RED_L   = RGBColor(0xFE, 0xE2, 0xE2)
C_GREEN   = RGBColor(0x05, 0x96, 0x69)
C_GREEN_L = RGBColor(0xD1, 0xFA, 0xE5)
C_TEAL    = RGBColor(0x08, 0x91, 0xB2)
C_TEAL_L  = RGBColor(0xCC, 0xF0, 0xF8)
C_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
C_PURPLE_L= RGBColor(0xED, 0xE9, 0xFE)
C_GREY    = RGBColor(0x64, 0x74, 0x8B)
C_GREY2   = RGBColor(0xF1, 0xF5, 0xF9)
C_NAVY    = RGBColor(0x1E, 0x3A, 0x8A)

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
blank = prs.slide_layouts[6]
W = 33.87
H = 19.05

# ── Layout ─────────────────────────────────────────────────────
HDR_H   = 1.22
CONT_Y  = 2.88
FOOT_Y  = H - 0.26    # 18.79
STRIP_Y = H - 1.22    # 17.83
CONT_H  = STRIP_Y - CONT_Y   # 14.95
CARD_W  = 10.55
CARD_XS = [1.1, 12.0, 22.9]

# ── Core Helpers ───────────────────────────────────────────────
def cm(v): return Cm(v)

def rect(slide, l, t, w, h, fill, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, cm(l), cm(t), cm(w), cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=10, bold=False,
        color=None, align=PP_ALIGN.LEFT, italic=False):
    if color is None: color = C_TEXT
    txb = slide.shapes.add_textbox(cm(l), cm(t), cm(w), cm(h))
    txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return txb

def txt_ml(slide, lines, l, t, w, h, size=10, bold=False,
           color=None, align=PP_ALIGN.LEFT):
    if color is None: color = C_TEXT
    txb = slide.shapes.add_textbox(cm(l), cm(t), cm(w), cm(h))
    txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return txb

def circle(slide, cx, cy, r, fill, text="", tsize=9, tc=None, bold=True):
    s = slide.shapes.add_shape(9, cm(cx-r), cm(cy-r), cm(r*2), cm(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    if text:
        if tc is None: tc = C_WHITE
        txb = slide.shapes.add_textbox(cm(cx-r), cm(cy-r), cm(r*2), cm(r*2))
        txb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        rr = txb.text_frame.paragraphs[0].add_run()
        rr.text = text; rr.font.size = Pt(tsize)
        rr.font.bold = bold; rr.font.color.rgb = tc
    return s

def circ_out(slide, cx, cy, r, stroke, fill=None, text="", tsize=9):
    """Outlined circle (ring)"""
    if fill is None: fill = C_SURF
    s = slide.shapes.add_shape(9, cm(cx-r), cm(cy-r), cm(r*2), cm(r*2))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = stroke; s.line.width = Pt(1.5)
    if text:
        txb = slide.shapes.add_textbox(cm(cx-r), cm(cy-r), cm(r*2), cm(r*2))
        txb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        rr = txb.text_frame.paragraphs[0].add_run()
        rr.text = text; rr.font.size = Pt(tsize)
        rr.font.bold = True; rr.font.color.rgb = stroke
    return s

def sdot(slide, cx, cy, size, fill):
    """Square dot bullet"""
    rect(slide, cx - size/2, cy - size/2, size, size, fill)

def pbar(slide, l, t, w, h, pct, color, bg=None):
    if bg is None: bg = C_DIV
    rect(slide, l, t, w, h, bg)
    if pct > 0:
        rect(slide, l, t, max(w * pct, 0.02), h, color)

def shadow_card(slide, l, t, w, h, fill=None, line=None, lw=0.3):
    """White card with drop shadow"""
    rect(slide, l + 0.08, t + 0.08, w, h, C_SHADOW)
    if fill is None: fill = C_SURF
    return rect(slide, l, t, w, h, fill, line=line, lw=lw)

# ── Layout Helpers ─────────────────────────────────────────────
def add_bg(slide):
    rect(slide, 0, 0, W, H, C_BG)

def add_header(slide, section_text=None, section_accent=None):
    if section_accent is None: section_accent = C_BLUE
    rect(slide, 0, 0, W, HDR_H, C_HDR)
    rect(slide, 0, HDR_H, W, 0.06, section_accent)
    # Logo mark: small colored block
    rect(slide, 0.80, 0.30, 0.07, 0.62, section_accent)
    txt(slide, "GAIA3D", 1.02, 0.36, 3.8, 0.52,
        size=10.5, bold=True, color=C_WHITE)
    if section_text:
        sw = max(len(section_text) * 0.27 + 1.2, 3.6)
        # pill-style badge
        rect(slide, 5.20, 0.24, sw, 0.74, section_accent)
        rect(slide, 5.20, 0.24, 0.18, 0.74, C_WHITE)   # white left stripe
        txt(slide, section_text, 5.44, 0.26, sw - 0.28, 0.70,
            size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, "광역권 AI 모빌리티 시뮬레이션 플랫폼  ·  3차년도 연구계획  ·  2026",
        W - 20.5, 0.38, 20.0, 0.50, size=9.5, color=C_TEXT3, align=PP_ALIGN.RIGHT)

def add_title(slide, title, subtitle="", slide_num=None, accent=None):
    if accent is None: accent = C_BLUE
    # Ghost slide number
    if slide_num is not None:
        txt(slide, f"{slide_num:02d}", W - 3.6, 1.18, 3.4, 1.20,
            size=42, bold=True, color=RGBColor(0xE2, 0xE8, 0xF0),
            align=PP_ALIGN.RIGHT)
    txt(slide, title, 1.1, 1.34, W - 6.2, 0.96, size=22, bold=True, color=C_TEXT)
    if subtitle:
        txt(slide, subtitle, 1.1, 2.30, W - 6.2, 0.46, size=9.5, color=C_TEXT3,
            italic=True)
    # Modern 3-segment underline
    rect(slide, 1.1,  2.82, 3.50, 0.07, accent)
    rect(slide, 4.72, 2.82, 1.20, 0.07, RGBColor(
        min(accent[0]+80, 255), min(accent[1]+80, 255), min(accent[2]+80, 255)))
    rect(slide, 6.04, 2.82, 0.55, 0.07, C_DIV)

def add_footer(slide):
    rect(slide, 0, FOOT_Y, W, H - FOOT_Y, C_HDR)
    rect(slide, 0, FOOT_Y, W, 0.06, C_BLUE)

def add_strip(slide, items, accent=None):
    if accent is None: accent = C_BLUE
    strip_h = FOOT_Y - STRIP_Y
    rect(slide, 0, STRIP_Y, W, strip_h, C_SURF, line=C_DIV, lw=0.3)
    rect(slide, 0, STRIP_Y, W, 0.06, accent)
    n = len(items); iw = (W - 2.2) / n
    for i, item in enumerate(items):
        if i > 0:
            rect(slide, 1.1 + i * iw, STRIP_Y + 0.20, 0.03, 0.58, C_DIV)
        # Square icon
        rect(slide, 1.10 + i * iw, STRIP_Y + 0.36, 0.16, 0.16, accent)
        txt(slide, item, 1.36 + i * iw, STRIP_Y + 0.22, iw - 0.54, 0.62,
            size=10.5, bold=True, color=C_TEXT2)

# ── Modern info_card (top-bar style) ───────────────────────────
def info_card(slide, l, t, w, h, accent, al, title, items,
              icon_text="", kpi=""):
    # Shadow
    rect(slide, l + 0.08, t + 0.08, w, h, C_SHADOW)
    # Card body
    rect(slide, l, t, w, h, C_SURF, line=C_BORDER, lw=0.3)
    # Top accent bar
    rect(slide, l, t, w, 0.44, accent)
    # Subtle tint just below bar
    rect(slide, l, t + 0.44, w, 0.18, al)
    # Icon circle in bar (top-right)
    if icon_text:
        circle(slide, l + w - 0.60, t + 0.22, 0.26, C_WHITE,
               icon_text, tsize=8, tc=accent)
    # Title
    title_w = w - 1.30 if icon_text else w - 0.48
    txt(slide, title, l + 0.28, t + 0.52, title_w, 0.72,
        size=11.5, bold=True, color=C_TEXT)
    # Separator
    rect(slide, l + 0.28, t + 1.28, w - 0.44, 0.03, C_DIV)
    # Items
    kpi_h = 0.50 if kpi else 0.0
    avail = h - 1.34 - kpi_h
    slot  = avail / len(items)
    for i, item in enumerate(items):
        iy   = t + 1.34 + i * slot
        mid  = iy + slot / 2
        # Square bullet
        sdot(slide, l + 0.44, mid, 0.09, accent)
        item_h = min(slot * 0.78, 1.05)
        txt(slide, item, l + 0.62, iy + (slot - item_h) / 2,
            w - 0.78, item_h, size=9.5, color=C_TEXT2)
    # KPI badge
    if kpi:
        by = t + h - 0.46
        rect(slide, l + 0.28, by, w - 0.44, 0.38, al, line=accent, lw=0.5)
        txt(slide, "▲ " + kpi, l + 0.28, by, w - 0.44, 0.38,
            size=8.5, bold=True, color=accent, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# Slide 1: Cover
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)

# Background gradient simulation
rect(slide, 0, 0, W, H, C_BG)
rect(slide, 0, 0, W * 0.60, H, C_SURF)   # left white panel

# Left accent stripe
rect(slide, 0, 0, 0.34, H, C_BLUE)
rect(slide, 0.34, 0, 0.06, H, C_BLUE_L)

# Right dark panel
rp_x = W * 0.60   # = 20.32
rect(slide, rp_x, 0, W - rp_x, H, C_BLUE_D)
rect(slide, rp_x, 0, 0.06, H, C_BLUE_M)

# ── IITP / year label ──
rect(slide, 0.70, 0.44, 3.8, 0.50, C_BLUE)
txt(slide, "IITP 정보통신기획평가원", 0.76, 0.46, 3.72, 0.46,
    size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── Main title ──
txt(slide, "광역권 도시교통", 0.70, 1.30, 18.4, 1.50,
    size=38, bold=True, color=C_TEXT)
txt(slide, "디지털트윈 시뮬레이션 플랫폼", 0.70, 2.84, 18.4, 1.30,
    size=28, bold=True, color=C_BLUE)

# Thick underline
rect(slide, 0.70, 4.24, 8.0, 0.10, C_BLUE)
rect(slide, 8.82, 4.24, 3.0, 0.10, C_BLUE_L)

txt(slide, "3차년도 연구계획  ·  2026년도", 0.70, 4.44, 18.0, 0.54,
    size=12, color=C_TEXT2)
txt(slide, "도시계획 지원 교통 시뮬레이션 구현 완성  |  LH 협력 실증 준비",
    0.70, 4.96, 18.0, 0.46, size=10.5, color=C_TEXT3)

# ── Tag pills ──
tags = [("신도시 시나리오", C_BLUE), ("도시정비 시나리오", C_TEAL),
        ("교통계획 평가", C_GREEN), ("분석 대시보드", C_AMBER), ("LH 실증 준비", C_PURPLE)]
px = 0.70
for tag, tc in tags:
    pw = len(tag) * 0.31 + 0.72
    rect(slide, px, 5.62, pw, 0.52, tc)
    txt(slide, tag, px + 0.12, 5.65, pw - 0.24, 0.46,
        size=9.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    px += pw + 0.18

# ── 4개년 진행 인포그래픽 ──
prog_y = 7.00
rect(slide, 0.70, prog_y + 0.30, 14.8, 0.06, C_DIV)   # base line
rect(slide, 0.70, prog_y + 0.30, 5.8,  0.06, C_GREY)   # completed

yrs = [
    ("1차\n2024", "기반 설계·데이터 모델",  "완료", C_GREY,  True,  False),
    ("2차\n2025", "프로토타입·핵심 기능",    "완료", C_GREY,  True,  False),
    ("3차\n2026", "구현 완성·시나리오",      "현재", C_BLUE,  False, True),
    ("4차\n2027", "LH 실증·서비스 전환",     "예정", C_BORDER,False, False),
]
for i, (lbl, desc, status, col, done, cur) in enumerate(yrs):
    cx = 0.70 + i * 4.92
    r = 0.40 if cur else 0.32
    if cur:
        # Outer glow ring
        circle(slide, cx, prog_y + 0.33, r + 0.14, C_BLUE_L)
    circle(slide, cx, prog_y + 0.33, r, col,
           "★" if cur else ("✓" if done else str(i+1)), tsize=10 if cur else 9)
    for j, line in enumerate(lbl.split("\n")):
        txt(slide, line, cx - 0.9, prog_y + 0.88 + j * 0.44, 1.8, 0.42,
            size=9 if not cur else 10, bold=cur, color=C_TEXT if cur else C_TEXT3,
            align=PP_ALIGN.CENTER)
    txt(slide, desc, cx - 1.1, prog_y + 1.80, 2.2, 0.60,
        size=8, color=C_TEXT2 if cur else C_TEXT3, align=PP_ALIGN.CENTER)
    bc = C_GREEN if done else (C_BLUE if cur else C_TEXT3)
    rect(slide, cx - 0.56, prog_y + 2.46, 1.12, 0.32, bc)
    txt(slide, status, cx - 0.56, prog_y + 2.47, 1.12, 0.30,
        size=8, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── 참여 기관 ──
rect(slide, 0.70, 10.20, 18.4, 0.04, C_DIV)
orgs = [("주관", "KAIST"), ("참여", "KOTI · Gaia3D · LaonRoad"), ("협력", "LH 한국토지주택공사")]
for i, (role, org) in enumerate(orgs):
    ox = 0.70 + i * 6.1
    txt(slide, role, ox, 10.30, 1.4, 0.34, size=8.5, color=C_TEXT3)
    txt(slide, org, ox + 1.1, 10.30, 4.8, 0.34, size=8.5, bold=True, color=C_TEXT2)

# ── 기술 스택 태그 ──
stacks = ["CesiumJS 3D", "OpenLayers 2D", "SUMO·VISSIM", "React + TS",
          "Spring Boot", "PostgreSQL/PostGIS"]
sx = 0.70
rect(slide, sx, 10.84, 18.4, 0.04, C_DIV)
for st in stacks:
    sw2 = len(st) * 0.23 + 0.54
    rect(slide, sx, 11.02, sw2, 0.44, C_SURF2, line=C_BORDER, lw=0.5)
    txt(slide, st, sx + 0.10, 11.04, sw2 - 0.14, 0.40, size=8.5, color=C_TEXT2)
    sx += sw2 + 0.14

# 연구 기간
txt(slide, "연구 기간  2024.07 ~ 2027.12  (총 4년, 정부출연금 기반 과제)",
    0.70, 11.64, 18.4, 0.40, size=9, color=C_TEXT3)

# ── 우측 패널 ──
rx = rp_x + 0.55
pw = W - rp_x - 0.82
txt(slide, "3차년도  구현 완성 단계", rx, 1.60, pw, 0.52,
    size=12.5, bold=True, color=C_BLUE_M)
txt_ml(slide, ["도시계획 지원", "교통 시뮬레이션"],
       rx, 2.18, pw, 1.58, size=20, bold=True, color=C_WHITE)
rect(slide, rx, 3.82, pw * 0.45, 0.06, C_BLUE_M)

stat_items = [
    ("3 / 4", "차년도 현재 진행",    C_BLUE_M,  0.75),
    ("Gaia3D","플랫폼 주관 기관",    RGBColor(0xFB,0xBF,0x24), 1.00),
    ("2027",  "플랫폼 완성 목표 연도", RGBColor(0x34,0xD3,0x99), 1.00),
]
card_sy = 4.24
card_sh = (FOOT_Y - card_sy - 0.22) / 3 - 0.11

for i, (val, lbl, col, pct) in enumerate(stat_items):
    sy = card_sy + i * (card_sh + 0.11)
    rect(slide, rx, sy, pw, card_sh, RGBColor(0x1E, 0x3A, 0x8A))
    rect(slide, rx, sy, 0.20, card_sh, col)
    # Top bar
    rect(slide, rx, sy, pw, 0.10, col)
    txt(slide, val, rx, sy + 0.36, pw, card_sh * 0.50,
        size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, lbl, rx, sy + card_sh * 0.56 + 0.20, pw, 0.44,
        size=9, color=C_BLUE_M, align=PP_ALIGN.CENTER)
    pbar(slide, rx + 0.34, sy + card_sh - 0.30, pw - 0.48, 0.14,
         pct, col, RGBColor(0x1E, 0x40, 0xAF))

rect(slide, 0, FOOT_Y, W, H - FOOT_Y, C_HDR)
rect(slide, 0, FOOT_Y, W, 0.06, C_BLUE)


# ═══════════════════════════════════════════════════════════════
# Slide 2: 전체 연구 계획 및 3차년도 위상
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide)
add_title(slide, "전체 연구 계획 및 3차년도 위상",
          "4개년 연구 로드맵  ·  현재 3차년도 구현 완성 단계  ·  2차년도 주요 성과 인계", slide_num=1)

# ── 4개년 타임라인 ──
tl_w = 7.56; tl_gap = (W - 2.2 - 4 * tl_w) / 3
tl_x0 = 1.1; tl_y = CONT_Y

year_info = [
    ("1차년도","2024","요구사항 분석\n기능 설계·DB 모델","완료",C_GREY,C_GREY2, False),
    ("2차년도","2025","프로토타입 구축\n핵심 기능 1차 구현","완료",C_GREY,C_GREY2, False),
    ("3차년도","2026","구현 완성\n시나리오·분석·대시보드","현재",C_BLUE,C_BLUE_L, True),
    ("4차년도","2027","LH 실증·고도화\n서비스 전환·안정화","예정",C_TEAL,RGBColor(0xF0,0xFD,0xFA), False),
]
th_list = [2.58, 2.58, 2.86, 2.58]

rect(slide, tl_x0 + 0.4, tl_y + 0.94, W - tl_x0 * 2 - 0.4, 0.06, C_DIV)

for i, (yn, yy, ytheme, yst, acc, al, cur) in enumerate(year_info):
    bx = tl_x0 + i * (tl_w + tl_gap)
    th = th_list[i]
    # Shadow + card
    rect(slide, bx + 0.06, tl_y + 0.06, tl_w, th, C_SHADOW)
    rect(slide, bx, tl_y, tl_w, th, C_SURF if cur else al, line=acc, lw=0.5 if cur else 0.3)
    # Top bar
    rect(slide, bx, tl_y, tl_w, 0.38, acc)
    # Year circle on bar
    circle(slide, bx + tl_w / 2, tl_y + 0.94, 0.36 if cur else 0.28,
           acc, "★" if cur else ("✓" if yst == "완료" else "○"), tsize=10 if cur else 9)
    # Status badge
    bc = C_GREEN if yst == "완료" else (C_BLUE if cur else C_GREY)
    rect(slide, bx + tl_w - 1.50, tl_y + 0.08, 1.42, 0.34, bc)
    txt(slide, yst, bx + tl_w - 1.50, tl_y + 0.08, 1.42, 0.34,
        size=8.5, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, yn, bx + 0.28, tl_y + 0.08, tl_w - 1.72, 0.32,
        size=10, bold=True, color=C_WHITE)
    txt(slide, yy, bx + 0.28, tl_y + 1.42, tl_w - 0.44, 0.38,
        size=10, color=acc if cur else C_TEXT3)
    for j, line in enumerate(ytheme.split("\n")):
        txt(slide, line, bx + 0.28, tl_y + 1.86 + j * 0.44, tl_w - 0.44, 0.42,
            size=9.5, color=C_TEXT2 if cur else C_TEXT3)
    if i < 3:
        ax = bx + tl_w + 0.04; ay = tl_y + 0.80
        arr = slide.shapes.add_shape(13, cm(ax), cm(ay), cm(tl_gap - 0.04), cm(0.38))
        arr.fill.solid()
        arr.fill.fore_color.rgb = C_BLUE if i == 1 else C_DIV
        arr.line.fill.background()

# ── 3차년도 5대 과제 ──
hl_y = tl_y + 3.12
rect(slide, 1.1, hl_y, W - 2.2, 0.04, C_BLUE)
txt(slide, "3차년도 5대 핵심 추진 과제", 1.1, hl_y + 0.10, 16, 0.40,
    size=9.5, bold=True, color=C_TEXT3)

tasks3 = [
    ("세부과제 01\n신도시/도시정비\n시나리오 구현",
     ["3기 신도시 6곳 교통계획 시뮬레이션",
      "재개발·재건축 교통영향평가 자동화",
      "노후계획도시 재정비 시나리오 구현",
      "A/B 비교 분석 대시보드 연동"],
     C_BLUE, C_BLUE_L, "01"),
    ("세부과제 02\n교통계획 평가\n시뮬레이션 완성",
     ["SUMO·VISSIM 속성 매핑 자동화",
      "OSM·국가교통DB Import 완성",
      "차종별 거동 모델 완전 구현",
      "REST API OpenAPI 3.0 표준화"],
     C_GREEN, C_GREEN_L, "02"),
    ("세부과제 03\n교통 분석\n고도화",
     ["HCM 6판 교차로 LOS 자동 산출",
      "구간 속도 프로파일 시계열 분석",
      "MOVES 2014b 탄소배출 추정",
      "A/B 시나리오 비교 뷰어 구현"],
     C_TEAL, C_TEAL_L, "03"),
    ("세부과제 04\n교통 분석\n대시보드 구현",
     ["실시간 KPI 6종 + 혼잡도 5단계",
      "시계열 차트 (교통량·속도) 완성",
      "차량 유형·속도 분포 차트 구현",
      "WebSocket 실시간 스트리밍"],
     C_AMBER, C_AMBER_L, "04"),
    ("세부과제 05\n플랫폼 완성\nLH 실증 준비",
     ["교통영향평가 PDF 자동 리포트",
      "3기 신도시 LH 실증 환경 구축",
      "CRDT 다중 사용자 협업 완성",
      "Docker + K8s 컨테이너 배포"],
     C_PURPLE, C_PURPLE_L, "05"),
]

tw = (W - 2.2 - 0.20 * 4) / 5
task_h = FOOT_Y - (hl_y + 0.60) - 0.06

for i, (ttl, descs, acc, al, num) in enumerate(tasks3):
    tx = 1.1 + i * (tw + 0.20)
    ty_t = hl_y + 0.60
    # Shadow + card
    rect(slide, tx + 0.05, ty_t + 0.05, tw, task_h, C_SHADOW)
    rect(slide, tx, ty_t, tw, task_h, C_SURF, line=C_BORDER, lw=0.3)
    rect(slide, tx, ty_t, tw, 0.34, acc)
    circle(slide, tx + tw / 2, ty_t + 0.93, 0.30, acc, num, tsize=9.5)
    for j, tline in enumerate(ttl.split("\n")):
        sz = 10.5 if j == 0 else 9.5; bld = (j == 0)
        col_t = C_WHITE if j == 0 else C_TEXT
        if j == 0:
            txt(slide, tline, tx + 0.16, ty_t + 0.04, tw - 0.32, 0.30,
                size=sz, bold=bld, color=C_WHITE)
        else:
            txt(slide, tline, tx + 0.16, ty_t + 1.30 + (j-1) * 0.48, tw - 0.32, 0.46,
                size=sz, bold=bld, color=C_TEXT)
    rect(slide, tx + 0.16, ty_t + 2.36, tw - 0.32, 0.03, C_DIV)
    slot_d = (task_h - 2.46) / len(descs)
    for k, d in enumerate(descs):
        dy = ty_t + 2.46 + k * slot_d
        sdot(slide, tx + 0.28, dy + slot_d / 2, 0.09, acc)
        txt(slide, d, tx + 0.46, dy + (slot_d - 0.84) / 2, tw - 0.60, 0.84,
            size=8.5, color=C_TEXT2)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 3: Contents
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide)
add_title(slide, "Contents",
          "3차년도 연구 추진 체계 및 세부 과제 구성", slide_num=2)

contents = [
    ("01","신도시/도시정비 시나리오 구현",
     "3기 신도시(남양주왕숙·하남교산·인천계양 등)·구도심 재개발·노후계획도시 재정비 3종 유스케이스 완전 구현",
     "시나리오 3종  ·  LH 협력 검증", C_BLUE, C_BLUE_L),
    ("02","교통계획 평가 시뮬레이션 완성",
     "SUMO·VISSIM 완전 연동, OSM·국가교통DB 자동 Import, 네트워크 편집 고도화, REST API 표준 완성",
     "Import 80% 단축  ·  OpenAPI 3.0", C_GREEN, C_GREEN_L),
    ("03","교통 분석 고도화",
     "HCM 6판 기준 교차로 LOS 자동 산출, 구간 속도 프로파일, MOVES 2014b 탄소배출 추정, A/B 비교",
     "3종 분석 자동화  ·  이상치 감지", C_TEAL, C_TEAL_L),
    ("04","교통 분석 대시보드 구현",
     "실시간 KPI 6종, 혼잡도 5단계, 시계열(교통량·속도) 차트, 차량·속도 분포 차트 완성",
     "완성도 100%  ·  비전문가 활용", C_AMBER, C_AMBER_L),
    ("05","플랫폼 완성 및 LH 실증 준비",
     "교통영향평가 PDF 자동 생성, 3기 신도시 LH 실증 환경, CRDT 협업, Docker·K8s 배포 자동화",
     "보고서 70% 단축  ·  실증 준비", C_PURPLE, C_PURPLE_L),
    ("06","추진 일정  —  분기별 로드맵",
     "Q1 기반 설계, Q2 핵심 기능 개발, Q3 기능 완성·연동, Q4 검증·배포 단계별 마일스톤 계획",
     "Q2 중간점검  ·  Q4 최종 완성", C_RED, C_RED_L),
    ("07","기대 효과 및 활용 방안",
     "신도시·행정·환경 3대 분야 활용, 보고서 70% 단축, LOS 자동화, 탄소 저감 정량화, 사회경제 파급효과",
     "정량 KPI 6종  ·  4개 활용 분야", C_BLUE, C_BLUE_L),
]

cw2 = (W - 2.2 - 0.22) / 2
rh = 2.72; rg = 0.16

for i, (num, ctitle, desc, kpi_txt, acc, al) in enumerate(contents):
    if i < 6:
        col = i % 2; row = i // 2
        lx = 1.1 + col * (cw2 + 0.22)
        ty = CONT_Y + row * (rh + rg)
        cw_use = cw2
    else:
        lx = 1.1; ty = CONT_Y + 3 * (rh + rg); cw_use = W - 2.2

    rect(slide, lx + 0.06, ty + 0.06, cw_use, rh, C_SHADOW)
    rect(slide, lx, ty, cw_use, rh, C_SURF, line=C_BORDER, lw=0.3)
    # Top accent bar
    rect(slide, lx, ty, cw_use, 0.22, acc)
    # Number badge on bar
    circle(slide, lx + 0.60, ty + 0.11, 0.22, C_WHITE, num, tsize=9, tc=acc)
    rect(slide, lx + 1.06, ty + 0.28, 0.03, rh - 0.44, C_DIV)
    txt(slide, ctitle, lx + 1.24, ty + 0.26, cw_use - 1.50, 0.56,
        size=12, bold=True, color=C_TEXT)
    txt(slide, desc, lx + 1.24, ty + 0.86, cw_use - 1.50, 1.14,
        size=8.5, color=C_TEXT2)
    rect(slide, lx + 1.24, ty + 2.06, cw_use - 1.50, 0.03, C_DIV)
    rect(slide, lx + 1.24, ty + 2.14, cw_use - 1.50, 0.42, al, line=acc, lw=0.4)
    txt(slide, "■  " + kpi_txt, lx + 1.34, ty + 2.16, cw_use - 1.70, 0.38,
        size=8.5, bold=True, color=acc)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 4: 세부과제 01 – 신도시/도시정비 시나리오
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide, section_text="세부과제 01", section_accent=C_BLUE)
add_title(slide, "신도시/도시정비 시나리오 구현",
          "3기 신도시 6곳 · 구도심 재개발·재건축 · 노후계획도시 재정비 — 3종 유스케이스 완전 구현",
          slide_num=3, accent=C_BLUE)

card_h = STRIP_Y - CONT_Y - 0.08

sc4 = [
    (C_BLUE, C_BLUE_L, "3기 신도시 개발 시나리오", "신", [
        "남양주왕숙·하남교산·인천계양·고양창릉·부천대장·과천 6개 신도시 시나리오 구현",
        "광역교통망 GTX-A·B·C / BRT / 도시철도 연계 수요 분석 및 환승 시뮬레이션",
        "시간대별 OD 행렬 기반 교통 수요 발생·배분·배정 모델 자동 적용",
        "신규 도로망 v/c 비율 · 평균 지체 · LOS 자동 산출 및 네트워크 시각화",
        "대중교통 노선·배차 간격 최적화 시뮬레이션 시나리오 구현",
        "교통대책 검토 결과 자동 PDF 보고서 생성 (LH 업무 연계)",
    ], "3기 신도시 6곳 직결"),
    (C_TEAL, C_TEAL_L, "구도심 재개발·재건축 시나리오", "재개발", [
        "재개발·재건축·리모델링 사업 구역 교통영향평가 프로세스 자동화",
        "사업 전/후 도로 용량 · 교차로 LOS · 대기행렬 Before/After 비교 분석",
        "보행 동선 · 자전거 인프라 · 교통약자 접근성 시뮬레이션",
        "주차 수요 분석 및 공영주차장 최적 위치 선정 알고리즘",
        "상업·업무·주거 복합 시설 유발 교통량 원단위법 추정 모델",
        "A/B 화면 분할 동기 재생 비교 대시보드 (KPI 6종 나란히 표시)",
    ], "A/B 비교 완전 지원"),
    (C_GREEN, C_GREEN_L, "노후계획도시 재정비 시나리오", "노후", [
        "노후계획도시 정비 특별법 기반 1기 신도시(분당·일산·평촌·산본·중동) 시나리오",
        "노후 교차로 LOS 개선 및 신호 재설계 최적화 시뮬레이션 구현",
        "간선버스 노선 재편 · 환승센터 신설 교통 개선 효과 정량 분석",
        "보행로 · 자전거도로 정비에 따른 수단 전환 효과 시뮬레이션",
        "단계별(1·2·3단계) 정비 효과 시계열 시각화 및 누적 KPI 표시",
        "LH · 지자체 협업 실증 데이터 수집 · 품질 관리 체계 구축",
    ], "LH 실증 연계"),
]

for i, (acc, al, ctitle, icon, items, kpi) in enumerate(sc4):
    info_card(slide, CARD_XS[i], CONT_Y, CARD_W, card_h,
              acc, al, ctitle, items, icon, kpi)

for i in range(2):
    ax = CARD_XS[i] + CARD_W + 0.08
    ay = CONT_Y + card_h / 2 - 0.22
    arr = slide.shapes.add_shape(13, cm(ax), cm(ay), cm(0.26), cm(0.44))
    arr.fill.solid(); arr.fill.fore_color.rgb = C_DIV
    arr.line.fill.background()

add_strip(slide, ["3종 시나리오 완전 구현", "LH 협력 교통대책 자동화", "A/B 비교 분석 지원"],
          accent=C_BLUE)
add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 5: 세부과제 02 – 시뮬레이션 완성
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide, section_text="세부과제 02", section_accent=C_GREEN)
add_title(slide, "교통계획 평가 시뮬레이션 완성",
          "SUMO·VISSIM 완전 연동  ·  다중 데이터 소스 Import  ·  모듈 간 REST API 표준화",
          slide_num=4, accent=C_GREEN)

pipe_h = 0.60
card_h5 = STRIP_Y - CONT_Y - pipe_h - 0.16

sc5 = [
    (C_BLUE, C_BLUE_L, "시뮬레이터 데이터 연동", "SIM", [
        "SUMO 0.32+ · VISSIM COM API 완전 연동 구현",
        "OSM · 국가교통DB · KTDB 자동 Import 엔진 완성",
        "WGS84 ↔ EPSG:5186 좌표계 자동 변환 처리",
        "네트워크 위상 정제 (dangling edge · isolated node 제거)",
        "Import 후 오류 자동 감지 및 사전 검증 체크 모듈",
        "시뮬레이션 파라미터 프리셋 저장·불러오기 기능",
    ], "Import 작업 80% 단축"),
    (C_TEAL, C_TEAL_L, "교통 시뮬레이션 구현", "교통", [
        "차로별 교통량 · 속도 실시간 CZML 3D 애니메이션",
        "차종별 거동 모델 (승용/화물/버스/자전거/보행자 5종)",
        "교통신호 현시 편집 및 SUMO 신호 연동 최적화",
        "OD 행렬 기반 교통수요 시뮬레이션 자동 적용",
        "대중교통 노선·정류장·배차 간격 편집 완전 구현",
        "돌발 상황(사고·공사 구간) 이벤트 시뮬레이션",
    ], "5종 차종 완전 구현"),
    (C_PURPLE, C_PURPLE_L, "모듈 간 인터페이스 · API", "API", [
        "REST API 표준 규격 OpenAPI 3.0 완전 문서화",
        "도시계획 레이어 데이터 → 시뮬레이션 입력 자동 변환",
        "수요예측 결과 → SUMO OD 행렬 연계 파이프라인",
        "시나리오별 결과 DB 저장 · 버전 관리 시스템",
        "WebSocket 실시간 시뮬레이션 진행 스트리밍 구현",
        "MSA 마이크로서비스 기반 통합 플랫폼 아키텍처",
    ], "3개 모듈 완전 통합"),
]

for i, (acc, al, ctitle, icon, items, kpi) in enumerate(sc5):
    info_card(slide, CARD_XS[i], CONT_Y, CARD_W, card_h5,
              acc, al, ctitle, items, icon, kpi)

# 데이터 파이프라인 인포그래픽
pipe_y = CONT_Y + card_h5 + 0.10
steps = [
    ("① 원본 데이터", "OSM / KTDB / VISSIM", C_BLUE_L, C_BLUE),
    ("② 속성 매핑", "자동 변환 엔진", C_TEAL_L, C_TEAL),
    ("③ 위상 정제", "오류 자동 제거", C_GREEN_L, C_GREEN),
    ("④ 시뮬레이션", "SUMO 실행·검증", C_AMBER_L, C_AMBER),
    ("⑤ 결과 저장", "DB 저장·API 제공", C_PURPLE_L, C_PURPLE),
]
pw2 = (W - 2.2) / len(steps)
for i, (st, sd, pc, pb) in enumerate(steps):
    px2 = 1.1 + i * pw2
    rect(slide, px2, pipe_y, pw2 - 0.06, pipe_h, pc, line=pb, lw=0.7)
    rect(slide, px2, pipe_y, pw2 - 0.06, 0.14, pb)
    txt(slide, st, px2 + 0.08, pipe_y + 0.16, pw2 - 0.20, 0.26,
        size=8, bold=True, color=pb)
    txt(slide, sd, px2 + 0.08, pipe_y + 0.36, pw2 - 0.20, 0.22,
        size=8, color=C_TEXT2)
    if i < len(steps) - 1:
        arr2 = slide.shapes.add_shape(13, cm(px2 + pw2 - 0.22),
                                       cm(pipe_y + (pipe_h - 0.28) / 2),
                                       cm(0.22), cm(0.28))
        arr2.fill.solid(); arr2.fill.fore_color.rgb = pb
        arr2.line.fill.background()

add_strip(slide, ["Import 80% 단축", "5종 차종 완전 구현", "3개 모듈 통합 완성"],
          accent=C_GREEN)
add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 6: 세부과제 03 – 교통 분석 고도화
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide, section_text="세부과제 03", section_accent=C_TEAL)
add_title(slide, "교통 분석 고도화",
          "HCM 6판 교차로 LOS  ·  구간 속도 프로파일  ·  MOVES 탄소배출  ·  A/B 시나리오 비교",
          slide_num=5, accent=C_TEAL)

# LOS 등급 인포그래픽
los_bar_h = 0.62
los_grades = [
    ("A", "매우 원활", "v/c ≤ 0.60",  C_GREEN),
    ("B", "원활",     "v/c ≤ 0.70",  RGBColor(0x65, 0xA3, 0x0D)),
    ("C", "보통",     "v/c ≤ 0.80",  C_AMBER),
    ("D", "지체",     "v/c ≤ 0.90",  RGBColor(0xEA, 0x58, 0x0C)),
    ("E", "심한지체", "v/c ≤ 1.00",  C_RED),
    ("F", "교통 마비","v/c > 1.00",  RGBColor(0x7F, 0x1D, 0x1D)),
]
gw = (W - 2.2) / 6
for j, (grade, glbl, vcratio, gc) in enumerate(los_grades):
    gx = 1.1 + j * gw
    rect(slide, gx, CONT_Y, gw - 0.05, los_bar_h, gc)
    txt(slide, grade, gx, CONT_Y + 0.04, gw - 0.05, 0.30,
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, glbl, gx, CONT_Y + 0.32, gw - 0.05, 0.22,
        size=7.5, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, vcratio, gx, CONT_Y + 0.50, gw - 0.05, 0.20,
        size=6.5, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
        italic=True)

card_h6 = STRIP_Y - CONT_Y - los_bar_h - 0.10
card_y6 = CONT_Y + los_bar_h + 0.06

sc6 = [
    (C_BLUE, C_BLUE_L, "교차로 성능 분석 (LOS)", "LOS", [
        "HCM 6판 기준 교차로 서비스 수준 A~F 자동 산출 엔진",
        "포화도(v/c ratio) · 평균 지체(초/대) · 대기행렬 실시간 계산",
        "현시별 포화교통류율 · 적정 현시율 자동 분석 및 제안",
        "차로별 LOS 히트맵 · 교차로 형태별(평면/회전) 비교",
        "HCM 기준 개선 전/후 LOS 변화 시뮬레이션 비교 뷰",
        "신호 최적화 알고리즘 기반 최적 현시 구성 자동 제안",
    ], "A~F 자동 판정 완성"),
    (C_GREEN, C_GREEN_L, "구간 속도 프로파일 분석", "SPD", [
        "5분 단위 시간대별 구간 속도 시계열 프로파일 생성",
        "Peak · Off-peak 자동 구분 및 혼잡 패턴 분석 리포트",
        "FCD 실측 데이터 vs 시뮬레이션 결과 정합성 비교 뷰",
        "링크별 속도 이상치(anomaly) 자동 감지 알림 시스템",
        "네트워크 전체 속도 히트맵 시각화 (레이어 오버레이)",
        "Recharts 기반 인터랙티브 속도 대시보드 연동",
    ], "이상치 자동 감지"),
    (C_AMBER, C_AMBER_L, "탄소배출 추정 · A/B 비교", "CO₂", [
        "MOVES 2014b 배출계수 기반 CO₂ 산출 엔진 완전 구현",
        "차종별 CO₂ · NOx · PM2.5 · HC 4종 배출량 정량화",
        "시나리오별 온실가스 감축 효과 및 탄소 크레딧 산정",
        "화면 분할 A/B 시나리오 동기 재생 비교 뷰어",
        "KPI 6종(교통량/속도/LOS/CO₂/NOx/PM2.5) 나란히 비교",
        "탄소 저감 최적 시나리오 자동 추천 알고리즘 구현",
    ], "3종 배출 자동 산출"),
]

for i, (acc, al, ctitle, icon, items, kpi) in enumerate(sc6):
    info_card(slide, CARD_XS[i], card_y6, CARD_W, card_h6,
              acc, al, ctitle, items, icon, kpi)

add_strip(slide, ["교차로 운영 효율 15% 향상", "탄소 저감 정량화", "최적 시나리오 자동 추천"],
          accent=C_TEAL)
add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 7: 교통 분석 대시보드 구현 현황
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide, section_text="세부과제 04", section_accent=C_AMBER)
add_title(slide, "교통 분석 대시보드 구현 현황 및 3차년도 고도화",
          "실시간 KPI 6종 · 혼잡도 5단계 · 시계열 차트 · A/B 비교 · 탄소배출 분석 완성",
          slide_num=6, accent=C_AMBER)

lc_w = 9.60
cx7  = 1.1 + lc_w + 0.26       # ≈ 10.96
cw7  = 12.60
rx7  = cx7 + cw7 + 0.26        # ≈ 23.82
rw7  = W - rx7 - 0.28          # ≈ 9.77

# ── 좌측: 구현 현황 ──
rect(slide, 1.00, CONT_Y, lc_w, 0.05, C_TEAL)
txt(slide, "2차년도 구현 완성 현황", 1.00, CONT_Y + 0.10, lc_w, 0.38,
    size=9.5, bold=True, color=C_TEXT3)

left_items = [
    ("KPI 카드 5종",     "활성 차량·도로 링크·편집 수·재생 속도·시뮬레이션 상태 실시간 표시", C_TEAL),
    ("실시간 진행률 바", "시뮬레이션 진행 % · 글로우 애니메이션 도트 실시간 표시", C_TEAL),
    ("탭 UI 전환",       "편집 이력 탭 ↔ 교통량 분석 탭 간 즉시 전환", C_BLUE),
    ("Live 현황 행",     "현재 시각 · 교통량 · 평균 속도 · 혼잡도 텍스트 실시간 표시", C_BLUE),
    ("혼잡도 5단계",     "심한 정체·정체·서행·원활·매우 원활 색상 인디케이터", C_BLUE),
    ("분석 KPI 6종",     "총 차량·여행 시간·평균 속도·피크량·이동 거리·CO₂ 자동 산출", C_GREEN),
    ("시계열 차트",      "교통량(막대)+속도(꺾은선)+현재 시각 수직선 실시간 표시", C_GREEN),
    ("분포 차트 2종",    "차량 유형 분포 + 속도 구간별 분포 막대 차트 완성", C_GREEN),
]

iw_l = (STRIP_Y - (CONT_Y + 0.56)) / len(left_items)
for i, (lt, ld, dc) in enumerate(left_items):
    sy = CONT_Y + 0.56 + i * iw_l
    sdot(slide, 1.12, sy + iw_l / 2, 0.11, dc)
    txt(slide, lt, 1.30, sy + (iw_l - 0.82) / 2, lc_w - 0.34, 0.40,
        size=9.5, bold=True, color=C_TEXT)
    txt(slide, ld, 1.30, sy + (iw_l - 0.82) / 2 + 0.44, lc_w - 0.34, 0.76,
        size=8, color=C_TEXT2)

# ── 중앙: 다크 UI 목업 ──
mock_h = STRIP_Y - CONT_Y
rect(slide, cx7 + 0.08, CONT_Y + 0.08, cw7, mock_h, C_SHADOW)
rect(slide, cx7, CONT_Y, cw7, mock_h,
     RGBColor(0x0C, 0x10, 0x18), line=RGBColor(0x2E, 0x3A, 0x55), lw=0.8)
# 타이틀 바
rect(slide, cx7, CONT_Y, cw7, 0.72, RGBColor(0x07, 0x09, 0x0F))
rect(slide, cx7, CONT_Y + 0.72, cw7, 0.05, C_BLUE_M)
txt(slide, "교통 시뮬레이션 대시보드  v3.0", cx7 + 0.24, CONT_Y + 0.16,
    cw7 * 0.58, 0.38, size=8.5, bold=True, color=RGBColor(0xE0, 0xE8, 0xF8))
txt(slide, "● 실행 중  07:28  |  62% 진행", cx7 + cw7 * 0.55, CONT_Y + 0.16,
    cw7 * 0.43, 0.38, size=7.5, color=RGBColor(0x10, 0xB9, 0x81),
    align=PP_ALIGN.RIGHT)

# KPI 5종
kpi7 = [("활성 차량","1,284",RGBColor(0x41,0x69,0xE1)),
        ("도로 링크", "892",  RGBColor(0x10,0xB9,0x81)),
        ("편집 수",   "7",    RGBColor(0xF5,0x9E,0x0B)),
        ("재생 속도", "3×",   RGBColor(0x8B,0x5C,0xF6)),
        ("상태",      "실행중",RGBColor(0x10,0xB9,0x81))]
kw7 = (cw7 - 0.20) / 5 - 0.04
for i, (kl, kv, kc) in enumerate(kpi7):
    kx7 = cx7 + 0.10 + i * (kw7 + 0.04)
    rect(slide, kx7, CONT_Y + 0.82, kw7, 1.20,
         RGBColor(0x14, 0x19, 0x29), line=RGBColor(0x23, 0x2E, 0x4A), lw=0.2)
    rect(slide, kx7, CONT_Y + 0.82, kw7, 0.10, kc)
    txt(slide, kl, kx7, CONT_Y + 0.96, kw7, 0.28,
        size=6.5, color=RGBColor(0x3A, 0x4A, 0x6A), align=PP_ALIGN.CENTER)
    txt(slide, kv, kx7, CONT_Y + 1.28, kw7, 0.52,
        size=11, bold=True, color=RGBColor(0xE0, 0xE8, 0xF8), align=PP_ALIGN.CENTER)

# 진행률 바
pb_y = CONT_Y + 2.12
rect(slide, cx7 + 0.10, pb_y, cw7 - 0.20, 0.18, RGBColor(0x0E, 0x14, 0x20))
rect(slide, cx7 + 0.10, pb_y, (cw7 - 0.20) * 0.62, 0.18, C_BLUE_M)

# 탭 바
rect(slide, cx7 + 0.10, pb_y + 0.26, cw7 - 0.20, 0.34, RGBColor(0x10, 0x14, 0x20))
for ti, (tlbl, tact) in enumerate([("교통량 분석", True), ("편집 이력", False), ("A/B 비교", False)]):
    tx_i = cx7 + 0.10 + ti * 3.2
    if tact:
        rect(slide, tx_i, pb_y + 0.26, 3.1, 0.34,
             RGBColor(0x16, 0x20, 0x3A), line=C_BLUE_M, lw=0.2)
    txt(slide, tlbl, tx_i + 0.10, pb_y + 0.28, 2.9, 0.30,
        size=7, color=C_WHITE if tact else RGBColor(0x3A, 0x4A, 0x6A))

# Live 행
rect(slide, cx7 + 0.10, pb_y + 0.68, cw7 - 0.20, 0.52,
     RGBColor(0x0E, 0x1E, 0x1E), line=RGBColor(0x06, 0xB6, 0xD4), lw=0.3)
circle(slide, cx7 + 0.26, pb_y + 0.94, 0.09, RGBColor(0x06, 0xB6, 0xD4))
txt(slide, "LIVE  07:28  |  교통량 1,284대  |  평균 속도 47.3 km/h  |  혼잡도: 서행",
    cx7 + 0.46, pb_y + 0.72, cw7 - 0.60, 0.38,
    size=7.5, color=RGBColor(0xE0, 0xE8, 0xF8))

# 혼잡도
seg7 = [RGBColor(0xEF,0x44,0x44), RGBColor(0xF5,0x9E,0x0B),
        RGBColor(0x10,0xB9,0x81), RGBColor(0x1A,0x22,0x30), RGBColor(0x1A,0x22,0x30)]
sw7 = (cw7 - 0.20) / len(seg7)
for si, sc in enumerate(seg7):
    rect(slide, cx7 + 0.10 + si * sw7, pb_y + 1.28, sw7 - 0.04, 0.16, sc)

# 분석 KPI
akpi = [("총 차량","4,201대"), ("여행 시간","8.2분"), ("평균 속도","51.3km/h"),
        ("피크 교통량","1,892대"), ("이동 거리","28.4km"), ("CO₂ 배출","3.4t/h")]
akw = (cw7 - 0.20) / 6 - 0.04
for i, (al2, av) in enumerate(akpi):
    akx = cx7 + 0.10 + i * (akw + 0.04)
    rect(slide, akx, pb_y + 1.52, akw, 1.06,
         RGBColor(0x14, 0x19, 0x29), line=RGBColor(0x23, 0x2E, 0x4A), lw=0.2)
    rect(slide, akx, pb_y + 1.52, akw, 0.08, C_BLUE_M)
    txt(slide, al2, akx, pb_y + 1.62, akw, 0.28,
        size=6, color=RGBColor(0x3A, 0x4A, 0x6A), align=PP_ALIGN.CENTER)
    txt(slide, av, akx, pb_y + 1.94, akw, 0.52,
        size=9.5, bold=True, color=RGBColor(0xE0, 0xE8, 0xF8), align=PP_ALIGN.CENTER)

# 차트 영역
ch_y = pb_y + 2.68
chart_h = STRIP_Y - ch_y - 0.06
main_cw = cw7 * 0.60 - 0.06
side_cw = cw7 - main_cw - 0.20 - 0.06

rect(slide, cx7 + 0.10, ch_y, main_cw, chart_h,
     RGBColor(0x14, 0x19, 0x29), line=RGBColor(0x23, 0x2E, 0x4A), lw=0.2)
txt(slide, "시간대별 교통량 · 속도 분석", cx7 + 0.22, ch_y + 0.08, main_cw - 0.20, 0.32,
    size=6.5, color=RGBColor(0x7A, 0x88, 0xAA))
bhs = [0.30, 0.46, 0.66, 0.86, 1.0, 0.80, 0.72, 0.92, 1.0, 0.58, 0.38, 0.30]
bar_w2 = (main_cw - 0.18) / len(bhs) - 0.02
for bi, bh in enumerate(bhs):
    bx7 = cx7 + 0.14 + bi * (bar_w2 + 0.02)
    bc = C_BLUE_M if bi < 8 else RGBColor(0x1C, 0x2E, 0x50)
    rect(slide, bx7, ch_y + chart_h - 0.12 - bh * (chart_h - 0.48),
         bar_w2, bh * (chart_h - 0.48), bc)
cur_x = cx7 + 0.14 + 8 * (bar_w2 + 0.02) - 0.02
rect(slide, cur_x, ch_y + 0.40, 0.04, chart_h - 0.52, RGBColor(0xEF, 0x44, 0x44))
txt(slide, "▲현재", cur_x - 0.22, ch_y + 0.08, 0.90, 0.26,
    size=5.5, color=RGBColor(0xEF, 0x44, 0x44))

side_x = cx7 + 0.10 + main_cw + 0.06
half_h = (chart_h - 0.06) / 2
for ci2, (ctitle2, vbars2, do_s) in enumerate([
    ("차량 유형 분포", [(C_BLUE_M,0.88),(C_AMBER,0.42),(C_GREEN,0.22),(C_PURPLE,0.12)], False),
    ("속도 구간 분포", [(C_RED,0.10),(RGBColor(0xF9,0x73,0x16),0.28),(C_AMBER,1.0),(C_GREEN,0.80),(C_BLUE_M,0.36)], True),
]):
    cy2 = ch_y + ci2 * (half_h + 0.06)
    rect(slide, side_x, cy2, side_cw, half_h,
         RGBColor(0x14, 0x19, 0x29), line=RGBColor(0x23, 0x2E, 0x4A), lw=0.2)
    txt(slide, ctitle2, side_x + 0.10, cy2 + 0.08, side_cw - 0.14, 0.28,
        size=6, color=RGBColor(0x7A, 0x88, 0xAA))
    vw2 = (side_cw - 0.20) / len(vbars2) - 0.04
    for vi2, (vc2, vh2) in enumerate(vbars2):
        vx2 = side_x + 0.10 + vi2 * (vw2 + 0.04)
        rect(slide, vx2, cy2 + half_h - 0.10 - vh2 * (half_h - 0.46),
             vw2, vh2 * (half_h - 0.46), vc2)

# ── 우측: 고도화 항목 (overflow 수정) ──
rect(slide, rx7, CONT_Y, rw7, 0.05, C_AMBER)
txt(slide, "3차년도 고도화 항목", rx7, CONT_Y + 0.10, rw7, 0.38,
    size=9.5, bold=True, color=C_TEXT3)

upg = [
    (C_BLUE, C_BLUE_L, "LOS 지표 연동",
     "HCM 6판 교차로 LOS A~F 자동 산출\n포화도·평균 지체(초/대) 실시간 표시\n신호 최적화 개선안 자동 제안"),
    (C_RED, C_RED_L, "A/B 시나리오 비교",
     "2개 시나리오 화면 분할 동기 재생\nKPI 6종 나란히 비교 대시보드\n최적 시나리오 자동 추천 기능"),
    (C_AMBER, C_AMBER_L, "탄소배출 상세 분석",
     "MOVES 2014b 차종별 CO₂·NOx·PM2.5\n온실가스 감축 효과 비교 분석\n탄소 저감 목표 달성률 시각화"),
    (C_PURPLE, C_PURPLE_L, "실시간 WebSocket 스트리밍",
     "WebSocket 기반 시뮬레이션 실시간 수신\n진행 중 시뮬레이션 모니터링 구현\n연결 끊김 자동 재연결 처리"),
    (C_GREEN, C_GREEN_L, "자동 리포트 생성",
     "KPI 기반 교통영향평가 PDF 자동 생성\n국토부 표준 보고서 포맷 완전 준수\nMP4 시뮬레이션 영상 자동 첨부"),
]

rs_y = CONT_Y + 0.52
gap_u = 0.08
uh = (STRIP_Y - rs_y - gap_u * (len(upg) - 1)) / len(upg)

for i, (acc, al, rt, rd) in enumerate(upg):
    uy = rs_y + i * (uh + gap_u)
    rect(slide, rx7 + 0.06, uy + 0.06, rw7, uh, C_SHADOW)
    rect(slide, rx7, uy, rw7, uh, C_SURF, line=C_BORDER, lw=0.3)
    rect(slide, rx7, uy, rw7, 0.32, acc)
    rect(slide, rx7, uy, 0.20, uh, acc)
    circle(slide, rx7 + 0.50, uy + 0.16, 0.18, C_WHITE, str(i + 1), tsize=8, tc=acc)
    txt(slide, rt, rx7 + 0.84, uy + 0.04, rw7 - 0.94, 0.26,
        size=9, bold=True, color=C_WHITE)
    rd_lines = rd.split("\n")
    txt_h = uh - 0.38
    txb = slide.shapes.add_textbox(cm(rx7 + 0.84), cm(uy + 0.36),
                                    cm(rw7 - 0.94), cm(txt_h))
    txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True
    for li, line in enumerate(rd_lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        r2 = p.add_run(); r2.text = line
        r2.font.size = Pt(8.5); r2.font.color.rgb = C_TEXT2

add_strip(slide, ["대시보드 완성도 100%", "의사결정 근거 시각화", "비전문가 활용 최적화"],
          accent=C_AMBER)
add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 8: 세부과제 05 – 플랫폼 완성 및 LH 실증
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide, section_text="세부과제 05", section_accent=C_PURPLE)
add_title(slide, "플랫폼 완성 및 LH 실증 준비",
          "교통영향평가 자동화 · LH 3기 신도시 실증 환경 · CRDT 협업 · Docker·K8s 배포",
          slide_num=7, accent=C_PURPLE)

card_h8 = STRIP_Y - CONT_Y - 0.08

sc8 = [
    (C_BLUE, C_BLUE_L, "자동 리포트 · 영상 내보내기", "RPT", [
        "시뮬레이션 KPI → 교통영향평가 보고서 자동 생성 엔진",
        "국토부 기준 HWP · PDF 표준 포맷 자동 작성 및 구성",
        "Cesium 뷰어 프레임 기반 시뮬레이션 화면 캡처 구현",
        "ffmpeg.wasm MP4 인코딩 — 브라우저 내 영상 변환",
        "개선안 비교표 · KPI 그래프 보고서 자동 삽입",
        "보고서 양식 템플릿 커스터마이징 및 기관별 적용",
    ], "보고서 작성 70% 단축"),
    (C_PURPLE, C_PURPLE_L, "LH 실증 환경 구축 준비", "LH", [
        "3기 신도시 6곳 교통계획 시나리오 사전 구현 완성",
        "노후계획도시 재정비 시뮬레이션 LH 검토 및 검증 완료",
        "LH 도시계획 업무 시스템 API 연계 인터페이스 개발",
        "LH 담당자 교육 프로그램 · 사용 매뉴얼 수립",
        "실증 데이터 수집 · 품질 관리 · 보안 체계 구축",
        "실증 성과 평가 지표 및 KPI 측정 기준 수립",
    ], "LH 실증 준비 완성"),
    (C_TEAL, C_TEAL_L, "협업 · 배포 환경 완성", "DEV", [
        "CRDT 알고리즘 기반 다중 사용자 실시간 동시 편집",
        "역할별 권한 관리 3단계 (관리자·편집자·뷰어)",
        "편집 이력 버전 관리 · 롤백 · 감사 로그 기능",
        "반응형 UI PC · 태블릿 해상도 완전 최적화",
        "Docker + Kubernetes 컨테이너 배포 자동화 구현",
        "GitHub Actions CI/CD 파이프라인 및 자동 테스트",
    ], "기관 간 협업 완성"),
]

for i, (acc, al, ctitle, icon, items, kpi) in enumerate(sc8):
    info_card(slide, CARD_XS[i], CONT_Y, CARD_W, card_h8,
              acc, al, ctitle, items, icon, kpi)

add_strip(slide, ["보고서 작성 70% 단축", "LH 실증 준비 완성", "기관 간 협업 구축"],
          accent=C_PURPLE)
add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 9: 추진 일정
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide)
add_title(slide, "3차년도 추진 일정  —  분기별 로드맵",
          "Q1 기반 설계 · Q2 핵심 기능 개발 · Q3 기능 완성·연동 · Q4 검증·배포 완성  ·  2026",
          slide_num=8)

qdata = [
    ("Q1  1 ~ 3월",  "기반 설계·착수",  C_BLUE,   C_BLUE_L),
    ("Q2  4 ~ 6월",  "핵심 기능 개발",  C_RED,    C_RED_L),
    ("Q3  7 ~ 9월",  "기능 완성·연동",  C_GREEN,  C_GREEN_L),
    ("Q4  10 ~ 12월","검증·배포·완성",  C_TEAL,   C_TEAL_L),
]

row_lbl_w = 2.24
qw = (W - row_lbl_w - 0.18 * 3) / 4
qx0 = row_lbl_w

for qi, (qtitle, qsub, qacc, qal) in enumerate(qdata):
    qx = qx0 + qi * (qw + 0.18)
    rect(slide, qx + 0.05, CONT_Y + 0.05, qw, 0.86, C_SHADOW)
    rect(slide, qx, CONT_Y, qw, 0.86, C_SURF, line=qacc, lw=0.6)
    rect(slide, qx, CONT_Y, qw, 0.24, qacc)
    txt(slide, qtitle, qx + 0.24, CONT_Y + 0.02, qw - 0.30, 0.24,
        size=9.5, bold=True, color=C_WHITE)
    txt(slide, qsub, qx + 0.24, CONT_Y + 0.28, qw - 0.30, 0.52,
        size=10, bold=True, color=C_TEXT)
    txt(slide, qsub, qx + 0.24, CONT_Y + 0.28, qw - 0.30, 0.52,
        size=10, bold=True, color=qacc)

rlabels = ["시나리오\n구현",  "시뮬레이션\n완성", "분석·\n대시보드", "플랫폼\n완성"]
rcolors = [C_BLUE, C_GREEN, C_TEAL, C_PURPLE]
rlights  = [C_BLUE_L, C_GREEN_L, C_TEAL_L, C_PURPLE_L]
n_rows   = len(rlabels)
rh9      = (FOOT_Y - (CONT_Y + 0.94) - 0.08 * (n_rows - 1)) / n_rows

roadmap = [
    [("신도시 시나리오 설계",   "유스케이스 요구사항 정의\n남양주왕숙·하남교산 설계",  "유스케이스 정의서 v1.0"),
     ("재개발·노후 시나리오",   "구도심 재개발 1차 구현\n노후도시 시나리오 착수",       "1차 시나리오 완성"),
     ("3종 시나리오 통합",      "3개 시나리오 통합 완성\nLH 검토 및 보완",              "통합 시나리오 완성"),
     ("LH 실증 검증·완성",      "LH 현장 적용 검증\n실증 결과 분석·평가",              "실증 완성")],
    [("SUMO 연동 설계·PoC",     "OSM Import 파서 프로토타입\n속성 매핑 아키텍처 설계", "파서 PoC 완성"),
     ("Import 엔진·편집 개발",  "OSM·KTDB Import 엔진 완성\n네트워크 편집 고도화 착수","Import 기능 완성"),
     ("편집·시뮬레이션 완성",   "네트워크·신호 편집 완성\n실데이터 통합 테스트",       "시뮬레이션 완성"),
     ("성능 최적화·검증",       "전체 성능 최적화 검증\nKPI 달성 최종 확인",           "검증 완료")],
    [("LOS 엔진 설계·PoC",      "HCM LOS 엔진 프로토타입\n교차로 분석 로직 설계",      "LOS PoC 완성"),
     ("속도·탄소배출 개발",     "속도 프로파일 대시보드\nMOVES 탄소배출 엔진 착수",    "분석 1차 완성"),
     ("A/B 비교·대시보드",      "탄소배출·A/B 비교뷰 완성\n대시보드 전체 통합",        "대시보드 완성"),
     ("분석 검증·성과 측정",    "HCM 검증·KPI 달성 확인\n최종 분석 성과 측정",         "분석 완성")],
    [("리포트·영상 설계",        "PDF 리포트 기본 구조 구현\nffmpeg 영상 캡처 PoC",     "기본 리포트 완성"),
     ("LH 실증 환경 착수",       "LH API 인터페이스 설계\n협업 기능·권한 관리 개발",   "LH 환경 착수"),
     ("협업·배포 환경 완성",     "CRDT 협업·권한 관리 완성\nDocker·K8s 배포 파이프라인","플랫폼 1차 완성"),
     ("배포·운영 전환",          "사용자 교육·인수·운영 전환\n최종 산출물 제출",        "배포 완료")],
]

for ri, (rl, rc, rli) in enumerate(zip(rlabels, rcolors, rlights)):
    ry = CONT_Y + 0.94 + ri * (rh9 + 0.08)

    # 행 레이블 (shadow 포함)
    rect(slide, 0.05, ry + 0.05, row_lbl_w, rh9, C_SHADOW)
    rect(slide, 0, ry, row_lbl_w, rh9, rli, line=rc, lw=0.4)
    rect(slide, 0, ry, 0.22, rh9, rc)
    rl_lines = rl.split("\n")
    rltxb = slide.shapes.add_textbox(cm(0.26), cm(ry + rh9 / 2 - 0.58),
                                      cm(row_lbl_w - 0.32), cm(1.18))
    rltxb.word_wrap = True; rltf = rltxb.text_frame; rltf.word_wrap = True
    for li, line in enumerate(rl_lines):
        p = rltf.paragraphs[0] if li == 0 else rltf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; r2 = p.add_run(); r2.text = line
        r2.font.size = Pt(8.5); r2.font.bold = True; r2.font.color.rgb = rc

    for ci in range(4):
        tcx = qx0 + ci * (qw + 0.18)
        ttl, tdet, ms = roadmap[ri][ci]
        rect(slide, tcx + 0.04, ry + 0.04, qw, rh9, C_SHADOW)
        rect(slide, tcx, ry, qw, rh9, C_SURF, line=C_BORDER, lw=0.3)
        rect(slide, tcx, ry, qw, 0.20, rli)
        rect(slide, tcx, ry, 0.18, rh9, rc)

        txt(slide, ttl, tcx + 0.28, ry + 0.10, qw - 0.38, 0.50,
            size=9.5, bold=True, color=C_TEXT)
        for di, dline in enumerate(tdet.split("\n")):
            txt(slide, dline, tcx + 0.28, ry + 0.66 + di * 0.44, qw - 0.38, 0.42,
                size=8, color=C_TEXT2)

        # 마일스톤 배지
        rect(slide, tcx + 0.28, ry + rh9 - 0.52, qw - 0.38, 0.40, rli, line=rc, lw=0.4)
        txt(slide, "◆  " + ms, tcx + 0.34, ry + rh9 - 0.50, qw - 0.44, 0.36,
            size=7.5, bold=True, color=rc)

        # 다이아몬드 마커
        dia = slide.shapes.add_shape(4,
            cm(tcx + qw - 0.34), cm(ry + 0.52), cm(0.28), cm(0.28))
        dia.fill.solid(); dia.fill.fore_color.rgb = rc
        dia.line.fill.background()

# 중간 점검선
mid_x = qx0 + (qw + 0.18) * 2 - 0.09
rect(slide, mid_x, CONT_Y + 0.06, 0.10, FOOT_Y - CONT_Y - 0.06,
     RGBColor(0xFE, 0xA3, 0xA3))
txt(slide, "중간 점검", mid_x - 0.24, CONT_Y + 0.10, 1.0, 0.34,
    size=7.5, bold=True, color=C_RED)

add_footer(slide)


# ═══════════════════════════════════════════════════════════════
# Slide 10: 기대 효과 및 활용 방안
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_bg(slide)
add_header(slide)
add_title(slide, "기대 효과 및 활용 방안",
          "정량 성과 목표 6종 · 신도시·행정·환경·분석 4대 활용 분야 · 사회경제적 파급효과",
          slide_num=9)

kpis = [
    ("3종",    "시나리오\n구현 완성",     C_BLUE,   C_BLUE_L,   1.00),
    ("LH",     "실증 준비\n완성",         C_GREEN,  C_GREEN_L,  1.00),
    ("15%↑",   "교차로 운영\n효율 향상",  C_TEAL,   C_TEAL_L,   0.85),
    ("100%",   "대시보드\n완성도",        C_AMBER,  C_AMBER_L,  1.00),
    ("70%↓",   "보고서 작성\n비용 절감",  C_PURPLE, C_PURPLE_L, 0.70),
    ("<500ms", "실시간\n응답 속도",       C_RED,    C_RED_L,    0.95),
]

kw = (W - 2.2) / len(kpis)
kcard_h = 3.22

for i, (kv, klbl, kacc, kal, kpct) in enumerate(kpis):
    kx = 1.1 + i * kw
    rect(slide, kx + 0.06, CONT_Y + 0.06, kw - 0.08, kcard_h, C_SHADOW)
    rect(slide, kx, CONT_Y, kw - 0.08, kcard_h, C_SURF, line=C_BORDER, lw=0.3)
    rect(slide, kx, CONT_Y, kw - 0.08, 0.30, kacc)
    txt(slide, kv, kx, CONT_Y + 0.34, kw - 0.08, 1.04,
        size=28, bold=True, color=kacc, align=PP_ALIGN.CENTER)
    kl_lines = klbl.split("\n")
    kltxb = slide.shapes.add_textbox(cm(kx), cm(CONT_Y + 1.42), cm(kw - 0.08), cm(0.92))
    kltxb.word_wrap = True; kltf = kltxb.text_frame; kltf.word_wrap = True
    for li, line in enumerate(kl_lines):
        p = kltf.paragraphs[0] if li == 0 else kltf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER; r2 = p.add_run(); r2.text = line
        r2.font.size = Pt(8.5); r2.font.color.rgb = C_TEXT2
    pbar(slide, kx + 0.26, CONT_Y + kcard_h - 0.46, kw - 0.50, 0.18, kpct, kacc, kal)
    txt(slide, f"{int(kpct*100)}%", kx + kw - 1.20, CONT_Y + kcard_h - 0.66, 1.00, 0.24,
        size=7.5, color=kacc, align=PP_ALIGN.RIGHT)

# 활용 방안 카드
uc_y = CONT_Y + kcard_h + 0.20
txt(slide, "주요 활용 분야", 1.1, uc_y, 14, 0.40, size=10, bold=True, color=C_TEXT3)
rect(slide, 1.1, uc_y + 0.44, W - 2.2, 0.04, C_DIV)

uc_data = [
    (C_BLUE, C_BLUE_L, "신도시·도시정비 활용", "신\n도시", [
        "3기 신도시 6곳 교통대책 수립 지원",
        "재개발·재건축 교통영향평가 자동화",
        "LH 업무 시스템 API 직접 연계",
        "노후 신도시 재정비 계획 지원",
    ]),
    (C_RED, C_RED_L, "행정·정책 의사결정", "행\n정", [
        "교통영향평가 자동화 보고서 제출",
        "지자체 교통 운영 의사결정 지원",
        "도로 신설·확장 사전 효과 검토",
        "대중교통 개편 효과 예측 분석",
    ]),
    (C_TEAL, C_TEAL_L, "분석·시각화 플랫폼", "분\n석", [
        "실시간 교통 모니터링 시스템",
        "시나리오 비교 정책 결정 도구",
        "KPI 리포트 공공 포털 공유",
        "교통 빅데이터 분석 플랫폼 연계",
    ]),
    (C_AMBER, C_AMBER_L, "환경·지속가능성 기여", "환\n경", [
        "탄소배출 저감 시나리오 비교",
        "친환경 교통 정책 효과 입증",
        "온실가스 감축 목표 달성 기여",
        "보행·자전거 전환 효과 정량화",
    ]),
]

uc_w  = (W - 2.2) / 4 - 0.18
uc_h  = STRIP_Y - uc_y - 0.58
uc_xs = [1.1 + i * ((W - 2.2) / 4) for i in range(4)]

for i, (acc, al, utitle, icon, uitems) in enumerate(uc_data):
    ux = uc_xs[i]; uy2 = uc_y + 0.60
    rect(slide, ux + 0.07, uy2 + 0.07, uc_w, uc_h, C_SHADOW)
    rect(slide, ux, uy2, uc_w, uc_h, C_SURF, line=C_BORDER, lw=0.3)
    rect(slide, ux, uy2, uc_w, 0.38, acc)
    circle(slide, ux + uc_w / 2, uy2 + 0.19, 0.22, C_WHITE,
           icon.split("\n")[0], tsize=8, tc=acc)
    txt(slide, utitle, ux + 0.18, uy2 + 0.44, uc_w - 0.24, 0.54,
        size=10.5, bold=True, color=C_TEXT)
    rect(slide, ux + 0.18, uy2 + 1.02, uc_w - 0.36, 0.03, C_DIV)
    slot_u = (uc_h - 1.10) / len(uitems)
    for j, uit in enumerate(uitems):
        iy = uy2 + 1.10 + j * slot_u
        sdot(slide, ux + 0.30, iy + slot_u / 2, 0.09, acc)
        txt(slide, uit, ux + 0.50, iy + (slot_u - 0.76) / 2, uc_w - 0.64, 0.76,
            size=9, color=C_TEXT2)

# 요약 배너
sb_y = STRIP_Y - 0.04
rect(slide, 1.1, sb_y, W - 2.2, FOOT_Y - sb_y - 0.04,
     C_BLUE_L, line=C_BLUE, lw=0.5)
rect(slide, 1.1, sb_y, 0.28, FOOT_Y - sb_y - 0.04, C_BLUE)
txt(slide, "3차년도 구현 완성  →  LH 실증  →  교통 디지털트윈 표준 플랫폼 확산",
    1.52, sb_y + 0.12, W - 3.0, 0.58, size=12.5, bold=True, color=C_TEXT)

add_strip(slide, ["신도시 교통계획 자동화", "LH 실증 완성", "탄소 저감 정량화"])
add_footer(slide)


# ── Save ──────────────────────────────────────────────────────
out = "/Users/hskim/Documents/repo/iitp/IITP_2026_3차년도_연구계획.pptx"
prs.save(out)
print(f"Saved: {out}")
