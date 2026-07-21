# -*- coding: utf-8 -*-
"""
IITP 네트워크 타일링/시뮬레이션 스트리밍 — 주요 변경 내역 및 개선 과제 PPT (상세 서술판)
각 주제를 "문제 상황 → 원인 → 해결 → 효과" 흐름으로 풀어서 설명.
(make_ppt_3rd.py 디자인 시스템 재사용: top-bar card style)
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color Palette ──────────────────────────────────────────────
C_BG      = RGBColor(0xF0, 0xF4, 0xF8)
C_HDR     = RGBColor(0x0D, 0x14, 0x26)
C_SURF    = RGBColor(0xFF, 0xFF, 0xFF)
C_SHADOW  = RGBColor(0xD8, 0xE3, 0xEE)
C_BORDER  = RGBColor(0xCB, 0xD5, 0xE1)
C_DIV     = RGBColor(0xE2, 0xE8, 0xF0)
C_TEXT    = RGBColor(0x0D, 0x14, 0x26)
C_TEXT2   = RGBColor(0x47, 0x55, 0x69)
C_TEXT3   = RGBColor(0x94, 0xA3, 0xB8)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BLUE    = RGBColor(0x1D, 0x4E, 0xD8)
C_BLUE_M  = RGBColor(0x3B, 0x82, 0xF6)
C_BLUE_L  = RGBColor(0xDB, 0xE9, 0xFE)
C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
C_AMBER_L = RGBColor(0xFE, 0xF3, 0xC7)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_RED_L   = RGBColor(0xFE, 0xE2, 0xE2)
C_GREEN   = RGBColor(0x05, 0x96, 0x69)
C_GREEN_L = RGBColor(0xD1, 0xFA, 0xE5)
C_TEAL    = RGBColor(0x08, 0x91, 0xB2)
C_TEAL_L  = RGBColor(0xCC, 0xF0, 0xF8)
C_PURPLE  = RGBColor(0x7C, 0x3A, 0xED)
C_PURPLE_L= RGBColor(0xED, 0xE9, 0xFE)
C_GREY2   = RGBColor(0xF1, 0xF5, 0xF9)

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
blank = prs.slide_layouts[6]
W, H = 33.87, 19.05
HDR_H = 1.22

def cm(v): return Cm(v)

def rect(slide, l, t, w, h, fill, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, cm(l), cm(t), cm(w), cm(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=10, bold=False,
        color=None, align=PP_ALIGN.LEFT):
    if color is None: color = C_TEXT
    txb = slide.shapes.add_textbox(cm(l), cm(t), cm(w), cm(h))
    txb.word_wrap = True; tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return txb

def sdot(slide, cx, cy, size, fill):
    rect(slide, cx - size/2, cy - size/2, size, size, fill)

def shadow_card(slide, l, t, w, h, fill=None, line=None, lw=0.3):
    rect(slide, l + 0.08, t + 0.08, w, h, C_SHADOW)
    if fill is None: fill = C_SURF
    return rect(slide, l, t, w, h, fill, line=line, lw=lw)

def add_bg(slide): rect(slide, 0, 0, W, H, C_BG)

def add_header(slide, section_text=None, accent=None):
    if accent is None: accent = C_BLUE
    rect(slide, 0, 0, W, HDR_H, C_HDR)
    rect(slide, 0, HDR_H, W, 0.06, accent)
    rect(slide, 0.80, 0.30, 0.07, 0.62, accent)
    txt(slide, "GAIA3D", 1.02, 0.36, 3.8, 0.52, size=10.5, bold=True, color=C_WHITE)
    if section_text:
        sw = max(len(section_text) * 0.30 + 1.4, 3.6)
        rect(slide, 5.20, 0.24, sw, 0.74, accent)
        rect(slide, 5.20, 0.24, 0.18, 0.74, C_WHITE)
        txt(slide, section_text, 5.44, 0.28, sw - 0.28, 0.66, size=11, bold=True, color=C_WHITE)

def add_title(slide, title, sub=None):
    txt(slide, title, 1.1, 1.50, 30, 0.95, size=18, bold=True)
    if sub:
        txt(slide, sub, 1.1, 2.28, 31, 0.60, size=10.5, color=C_TEXT2)

def card_head(slide, l, t, w, label, accent):
    rect(slide, l, t, w, 0.62, accent)
    txt(slide, label, l + 0.25, t + 0.07, w - 0.4, 0.5, size=11, bold=True, color=C_WHITE)

def labeled_para(slide, l, t, w, label, body, label_color, body_size=9.5, body_h=1.4):
    """라벨(굵게) + 본문 문단"""
    txt(slide, label, l, t, w, 0.45, size=10, bold=True, color=label_color)
    txt(slide, body, l, t + 0.44, w, body_h, size=body_size, color=C_TEXT2)

def story_block(slide, l, t, w, tag, tag_color, tag_bg, title, body, body_h=1.15):
    """[문제]/[원인]/[해결]/[효과] 태그가 붙은 서술 블록. 반환값: 다음 y"""
    rect(slide, l, t + 0.03, 1.35, 0.5, tag_bg)
    txt(slide, tag, l, t + 0.06, 1.35, 0.45, size=9, bold=True, color=tag_color, align=PP_ALIGN.CENTER)
    txt(slide, title, l + 1.55, t + 0.03, w - 1.6, 0.5, size=10.5, bold=True)
    txt(slide, body, l + 1.55, t + 0.50, w - 1.6, body_h, size=9.5, color=C_TEXT2)
    return t + 0.50 + body_h + 0.18

def metric(slide, l, t, w, big, small, accent, h=2.15):
    shadow_card(slide, l, t, w, h)
    rect(slide, l, t, 0.10, h, accent)
    txt(slide, big, l + 0.30, t + 0.28, w - 0.5, 0.85, size=19, bold=True, color=accent)
    txt(slide, small, l + 0.30, t + 1.18, w - 0.5, h - 1.2, size=9, color=C_TEXT2)

TAG_PROB = ("문제",  C_RED,    C_RED_L)
TAG_CAUSE= ("원인",  C_AMBER,  C_AMBER_L)
TAG_FIX  = ("해결",  C_BLUE,   C_BLUE_L)
TAG_EFF  = ("효과",  C_GREEN,  C_GREEN_L)

# ════════════════════════════════════════════════════════════════
# 1 — 표지
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s)
rect(s, 0, 0, W, H, C_HDR)
rect(s, 0, 12.3, W, 0.10, C_BLUE_M)
txt(s, "광역권 모빌리티 시뮬레이션 플랫폼", 2.2, 4.9, 30, 1.0, size=15, color=C_BLUE_L)
txt(s, "네트워크 타일링 · 시뮬레이션 스트리밍", 2.2, 5.9, 30, 1.6, size=30, bold=True, color=C_WHITE)
txt(s, "주요 변경 내역 및 개선 과제 (상세)", 2.2, 7.6, 30, 1.0, size=20, bold=True, color=C_WHITE)
txt(s, "실증 데이터: 대전·세종 전역 KTDB 표준노드링크 — 링크 62,434개 / 노드 47,181개 / 더미 차량 186,110대",
    2.2, 9.2, 30, 0.7, size=12, color=C_TEXT3)
txt(s, "GAIA3D  ·  IITP 광역권 도시 차세대 AI 융합 모빌리티 과제", 2.2, 13.0, 30, 0.7, size=11, color=C_TEXT3)

# ════════════════════════════════════════════════════════════════
# 2 — 배경: 왜 이 작업이 필요했나
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "00 배경", C_TEXT3)
add_title(s, "왜 이 작업이 필요했나 — 전체 일괄 로드 구조의 물리적 한계",
          "기존 구조는 시나리오를 열 때 네트워크와 시뮬레이션 데이터 전체를 브라우저에 내려받아 한 번에 그리는 방식")

card_head(s, 1.1, 3.35, 15.4, "기존 구조 (전체 일괄 로드)", C_RED)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
y = 4.35
y = story_block(s, 1.5, y, 14.6, *TAG_PROB, "시나리오를 열면 화면이 멈추거나 아예 뜨지 않음",
    "대전·세종 네트워크는 XML 106MB. 브라우저가 이걸 통째로 다운로드→파싱→3D 오브젝트로 빌드하는 동안 "
    "화면이 수십 초간 멈추고, 느린 환경에서는 소켓 타임아웃으로 아예 빈 화면이 됨.", body_h=1.35)
y = story_block(s, 1.5, y, 14.6, *TAG_PROB, "메모리 한계 — 대도시 1개가 상한",
    "실측 결과 전체-로드 방식은 브라우저 8GB 힙에서 약 20만 링크가 한계. 수도권 환산 시 "
    "네트워크만 ~1GB, 차량 시뮬레이션은 수 GB 이상이라 광역권 확장이 원천 불가능.", body_h=1.35)
y = story_block(s, 1.5, y, 14.6, *TAG_PROB, "시뮬레이션 결과도 같은 구조",
    "차량 궤적을 CZML(전체 차량·전체 시간)로 한 번에 반환. 대전·세종 규모(18.6만 대)는 "
    "생성 단계에서 서버 메모리 초과(OOM)로 크래시 — 생성 자체가 불가능했음.", body_h=1.35)

card_head(s, 17.3, 3.35, 15.4, "전환 방향 (viewport 스트리밍)", C_GREEN)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
txt(s, "핵심 아이디어: “지금 화면에 보이는 것만, 보이는 수준(LOD)으로”",
    17.7, 4.30, 14.6, 0.6, size=11, bold=True, color=C_GREEN)
bl = [
    ("공간을 자른다 (타일)", "지도를 경위도 격자로 나눠, 카메라 화면과 겹치는 타일만 서버에서 받아온다. 화면을 벗어난 타일은 메모리에서 회수(LRU)."),
    ("디테일을 자른다 (LOD)", "멀리서는 고속도로 골격만, 가까이서는 차선·교차로 구성요소까지 — 카메라 거리(픽셀 해상도)에 따라 5단계로 응답 내용을 다르게."),
    ("시간을 자른다 (시간창)", "차량 궤적은 재생 중인 시각 주변 몇 분 치만 로드하고, 창이 소진되기 전에 다음 창을 미리 받아 끊김 없이 이어재생."),
    ("서버가 자른다 (인덱스)", "자르는 계산은 서버의 공간 인덱스(SQLite RTree)가 담당 — 클라이언트는 조각만 받으므로 데이터 총량과 무관하게 동작."),
]
yy = 5.05
for head, body in bl:
    sdot(s, 17.85, yy + 0.20, 0.13, C_GREEN)
    txt(s, head, 18.10, yy, 14.2, 0.45, size=10.5, bold=True)
    txt(s, body, 18.10, yy + 0.44, 14.1, 1.15, size=9.5, color=C_TEXT2)
    yy += 1.90

# ════════════════════════════════════════════════════════════════
# 3 — 성과 요약 (Before/After 표)
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "01 핵심 성과", C_BLUE)
add_title(s, "핵심 성과 요약 — Before / After",
          "아래 수치는 모두 대전·세종 실데이터로 브라우저(E2E)·API에서 직접 측정한 값")

rows = [
    ("네트워크 초기 로드", "106MB 전체 다운로드+파싱 → 수십 초 멈춤,\n타임아웃 시 빈 화면", "viewport 타일만 로드 — 진입 후 ~17초에 지도 표시,\n네트워크 위치로 카메라 자동 이동", C_BLUE),
    ("시뮬레이션 차량", "18.6만 대 생성 시도 → 서버 OOM 크래시\n(이벤트 4,830만 건을 메모리에 전부 적재)", "스트리밍 생성 70초 완료. 재생은 화면·시간창 단위\n(화면당 최대 1,500대, 초과 시 히트맵 자동 전환)", C_TEAL),
    ("브라우저 메모리", "JS Heap ~1GB 이상, 프레임 한 자리까지 하락", "220MB 수준 — 화면 밖 데이터를 지우는\n거리 기반 생명주기로 총량 자체를 제한", C_GREEN),
    ("서버 디스크", "임시 DB 사본이 요청마다 생성·미삭제\n→ 70GB 누적 (디스크 풀 실사고)", "시나리오당 사본 1개 캐시 재사용 + 데이터 변경 시\n무효화 연동 — 누수 원천 차단", C_PURPLE),
]
y0 = 3.3
colX = [1.1, 8.0, 17.6]
txt(s, "항목", colX[0]+0.2, y0, 5, 0.5, size=11, bold=True, color=C_TEXT2)
txt(s, "이전 (Before)", colX[1]+0.2, y0, 8, 0.5, size=11, bold=True, color=C_RED)
txt(s, "현재 (After)", colX[2]+0.2, y0, 8, 0.5, size=11, bold=True, color=C_GREEN)
y = y0 + 0.62
for name, before, after, ac in rows:
    shadow_card(s, 1.1, y, 31.6, 2.62)
    rect(s, 1.1, y, 0.10, 2.62, ac)
    txt(s, name, colX[0]+0.35, y+0.90, 6.3, 0.9, size=12, bold=True)
    for i, line in enumerate(before.split("\n")):
        txt(s, line, colX[1]+0.2, y+0.62+i*0.62, 9.2, 0.65, size=9.5, color=C_TEXT2)
    txt(s, "→", 16.60, y+0.95, 1.0, 0.6, size=14, bold=True, color=C_TEXT3)
    for i, line in enumerate(after.split("\n")):
        txt(s, line, colX[2]+0.2, y+0.62+i*0.62, 15.0, 0.65, size=9.5, color=C_TEXT)
    y += 2.90

# ════════════════════════════════════════════════════════════════
# 4 — 네트워크 타일링 ① 동작 원리
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "02 네트워크 타일링", C_BLUE)
add_title(s, "네트워크 타일링 ① — 어떻게 동작하나",
          "임포트할 때 공간 인덱스를 한 번 만들어 두고, 이후엔 카메라가 움직일 때마다 화면 범위(bbox)만 조회")

card_head(s, 1.1, 3.35, 15.4, "데이터 흐름", C_BLUE)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
steps = [
    ("1. 임포트 (1회)", "KTDB 표준노드링크를 network.xml로 변환하면서, 링크·노드를 SQLite에 넣고 RTree 공간 인덱스를 빌드. "
     "파일은 서버 디스크에 영속 저장 → 서버를 재시작해도 재임포트 불필요."),
    ("2. 시나리오 진입", "클라이언트는 전체 데이터를 받지 않는다. 대신 extent API로 네트워크 전체 범위만 조회해 "
     "카메라를 그 위치로 이동 (시나리오 등록 좌표가 실제 데이터 위치와 다른 사례가 있어 등록 좌표는 신뢰하지 않음)."),
    ("3. 카메라 정착 시", "화면 중앙의 지면 거리로 픽셀 해상도(m/px)를 계산해 LOD 단계를 정하고, "
     "화면이 덮는 bbox와 겹치는 타일만 요청. 이미 받은 타일은 캐시에서 재사용."),
    ("4. 화면 이탈 시", "안 보이게 된 타일은 LRU로 회수(최대 64타일 유지). 진행 중이던 요청도 즉시 취소(Abort)해 "
     "빠른 줌/팬 중 중간 단계 요청이 쌓여 로딩이 밀리는 것을 방지."),
]
yy = 4.35
for head, body in steps:
    txt(s, head, 1.5, yy, 14.6, 0.45, size=10.5, bold=True, color=C_BLUE)
    txt(s, body, 1.5, yy + 0.44, 14.6, 1.35, size=9.5, color=C_TEXT2)
    yy += 2.05

card_head(s, 17.3, 3.35, 15.4, "LOD 5단계 — 거리에 따라 다른 응답", C_TEAL)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
lod_rows = [
    ("macro",    "~26km+",     "고속도로급 골격만 (속도 90+ 또는 6차선+).\n전국/광역 조망에서 큰 구조가 보이도록", C_PURPLE),
    ("overview", "8~26km",     "간선도로 중심선. 차선수에 비례한 굵기로\n도로 위계가 눈에 들어오게", C_BLUE),
    ("mid",      "1.5~8km",    "집산도로까지 포함한 중심선 — 줌인할수록 조밀", C_TEAL),
    ("near",     "310m~1.5km", "도로 폭 폴리곤(차선 데이터 제외한 경량).\n+3배 넓은 중심선 베이스로 타일 경계 끊김 방지", C_AMBER),
    ("detail",   "~310m",      "차선 구분선·중앙선 + 노드/포트/커넥션.\n응답에서 셀·구간 등 렌더 미사용 데이터 제거(수 배 축소)", C_GREEN),
]
ly = 4.35
for name, rng, desc, ac in lod_rows:
    rect(s, 17.7, ly, 2.5, 0.60, ac)
    txt(s, name, 17.7, ly+0.06, 2.5, 0.5, size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s, rng, 20.45, ly+0.06, 2.9, 0.5, size=9.5, bold=True, color=C_TEXT2)
    for i, line in enumerate(desc.split("\n")):
        txt(s, line, 23.4, ly+0.02+i*0.44, 9.2, 0.5, size=9, color=C_TEXT2)
    ly += 1.62

# ════════════════════════════════════════════════════════════════
# 5 — 네트워크 타일링 ② 안정화 여정
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "02 네트워크 타일링", C_BLUE)
add_title(s, "네트워크 타일링 ② — “네트워크가 사라진다” 를 잡기까지",
          "실사용 중 간헐적으로 도로가 통째로 사라지던 증상 — 서로 겹친 4개의 독립 원인을 재현·수정")

card_l = [1.1, 17.3]
card_head(s, card_l[0], 3.30, 15.4, "원인 1·2 — 지워 놓고 다시 못 그리는 레이스", C_RED)
shadow_card(s, card_l[0], 3.92, 15.4, 4.35)
y = 4.28
y = story_block(s, 1.5, y, 14.6, *TAG_CAUSE, "load() 재진입 시 전면 삭제",
    "2D 지도가 데이터를 동기화할 때마다 3D 초기화가 재진입해 모든 도로를 지웠는데, 타일 캐시는 "
    "“이미 로드됨” 상태라 다시 그리지 않음 → 영구 소실. 모드 전환 시에만 정리하도록 수정.", body_h=1.10)
y = story_block(s, 1.5, y, 14.6, *TAG_CAUSE, "지형 미로딩 시 조회 포기",
    "줌 직후 지형 타일이 아직 없으면 화면 중심 계산이 실패하는데 재시도가 없어, 카메라가 멈춰 있으면 "
    "타일 요청이 영영 시작되지 않음 → 타원체 폴백 + 1초 재시도 추가.", body_h=1.10)

card_head(s, card_l[0], 8.60, 15.4, "원인 3·4 — 렌더 루프와 LOD 경계", C_RED)
shadow_card(s, card_l[0], 9.22, 15.4, 4.35)
y = 9.58
y = story_block(s, 1.5, y, 14.6, *TAG_CAUSE, "요청 시 렌더 모드의 함정",
    "성능을 위해 “필요할 때만 프레임을 그리는” 모드 사용 중 — 카메라가 멈춘 뒤 도착한 도로는 "
    "프레임이 돌지 않아 준비가 영영 끝나지 않음. 준비될 때까지 렌더를 요청하는 펌프 추가(상한 없음).", body_h=1.10)
y = story_block(s, 1.5, y, 14.6, *TAG_CAUSE, "LOD 경계에서의 왕복 진동",
    "near/detail 경계 고도에서 계산이 순간 요동하면 전체 타일을 지우고 다시 받기를 반복. "
    "4초 이상 안정적으로 줌아웃일 때만 회수하는 히스테리시스 적용.", body_h=1.10)

card_head(s, card_l[1], 3.30, 15.4, "검증 방법", C_GREEN)
shadow_card(s, card_l[1], 3.92, 15.4, 9.65)
txt(s, "Playwright 브라우저 자동화로 실제 앱을 구동해 재현 → 수정 → 재검증을 반복",
    17.7, 4.25, 14.6, 0.6, size=10.5, bold=True, color=C_GREEN)
vitems = [
    ("증상 재현 스크립트", "경계 고도 왕복(260↔500m ×6), detail 연속 팬 ×6, 줌아웃 후 즉시 복귀, 90초 방치 — "
     "각 시나리오에서 화면의 도로 프리미티브 수를 직접 계측"),
    ("수정 전", "30초 방치 시 도로 수 2→0 (소실 재현), 250m 순간이동 후 타일 요청 0건"),
    ("수정 후", "전 시나리오에서 도로 수 유지(29개 안정), 요청 정상 발화, 페이지 에러 0건"),
    ("부수 발견", "줌 애니메이션 중 중간 배율의 타일 요청 수백 개가 브라우저 동시연결(6개)을 고갈시켜 "
     "“무한 로딩”처럼 보이던 문제 — 화면 밖 요청 즉시 취소로 해소"),
]
yy = 5.0
for head, body in vitems:
    sdot(s, 17.85, yy + 0.20, 0.13, C_GREEN)
    txt(s, head, 18.10, yy, 14.2, 0.45, size=10.5, bold=True)
    txt(s, body, 18.10, yy + 0.44, 14.1, 1.30, size=9.5, color=C_TEXT2)
    yy += 1.95

# ════════════════════════════════════════════════════════════════
# 6 — 3D 렌더링 ① 지형 정합
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "03 3D 렌더링", C_TEAL)
add_title(s, "3D 렌더링 ① — 도로가 땅속에 있거나, 땅에 묻히거나",
          "높이를 직접 지정하던 방식의 구조적 딜레마를 “지형 표면에 붙이는” 방식으로 전환")

card_head(s, 1.1, 3.35, 15.4, "문제와 원인", C_RED)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
y = 4.35
y = story_block(s, 1.5, y, 14.6, *TAG_PROB, "도로가 지형을 투과해 땅속에 보임",
    "도로를 그릴 높이를 “지형 고도 샘플링 평균 + 몇 cm”로 직접 계산했는데, 지형 서버 응답 실패나 "
    "고도 불일치가 있으면 도로가 고도 0(해수면)에 그려짐. 대전은 지대가 높아 도로가 수십 m 땅속에 위치.", body_h=1.55)
y = story_block(s, 1.5, y, 14.6, *TAG_CAUSE, "깊이 검사의 딜레마",
    "깊이 검사를 끄면 땅속 도로가 지형을 뚫고 비쳐 보이고(투과), 켜면 지형보다 낮은 도로가 완전히 "
    "가려져 안 보임. 높이를 직접 지정하는 한 어느 쪽을 선택해도 문제가 남는 구조.", body_h=1.55)
story_block(s, 1.5, y, 14.6, *TAG_FIX, "GroundPrimitive (지형 클램프) 전환",
    "높이 계산을 아예 버리고 Cesium이 도로를 지형 표면에 “드레이프”하도록 변경. 지형이 늦게 로드되거나 "
    "고도가 달라도 항상 표면 위. 부수 효과로 지형 샘플링 대기가 사라져 타일 표시가 즉시화.", body_h=1.55)

card_head(s, 17.3, 3.35, 15.4, "함께 정리한 정합 이슈", C_TEAL)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
vitems = [
    ("3D Tiles와의 공존", "도로를 지형에만 붙이면 3D 건물/정밀지형 타일셋을 켰을 때 도로가 그 밑에 가려짐 — "
     "지형+타일셋 양쪽 표면에 드레이프(BOTH)로 통일"),
    ("상·하행 도로 분리", "KTDB는 상행과 하행이 같은 중심선 좌표를 공유하는 별도 링크 — 그대로 그리면 도로 두 장이 "
     "완전히 겹침. 진행방향 우측으로 반 폭 이동시켜 실도로처럼 중앙선 양쪽으로 분리"),
    ("좌표 정제", "클러스터 병합 과정에서 생긴 급반전 좌표(150° 이상 꺾임)가 차선 계산을 무너뜨려 "
     "구분선이 직각으로 꺾이던 것 — 좌표 정제 단계에서 필터링"),
    ("카메라 자동 이동", "시나리오 진입 시 서버 extent API(네트워크 실제 범위)로 카메라 이동 — "
     "등록 좌표(부천)와 데이터(대전)가 다른 실사례 대응"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 17.85, yy + 0.20, 0.13, C_TEAL)
    txt(s, head, 18.10, yy, 14.2, 0.45, size=10.5, bold=True)
    txt(s, body, 18.10, yy + 0.44, 14.1, 1.30, size=9.5, color=C_TEXT2)
    yy += 2.05

# ════════════════════════════════════════════════════════════════
# 7 — 3D 렌더링 ② 시점 아티팩트와 성능
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "03 3D 렌더링", C_TEAL)
add_title(s, "3D 렌더링 ② — “카메라를 돌리면 무늬가 돌아간다”",
          "지형 드레이프 렌더의 시점 의존 아티팩트 — 완화가 아니라 원인 요소를 제거하는 방식으로 해결")

card_head(s, 1.1, 3.35, 15.4, "시점 의존 아티팩트", C_RED)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
y = 4.35
y = story_block(s, 1.5, y, 14.6, *TAG_PROB, "레인 회색면이 카메라 각도 따라 뒤틀림",
    "차선마다 깔던 회색 교차 음영(레인 채움면)이 카메라를 돌리면 함께 돌아가는 것처럼 보임. "
    "지형 드레이프는 3차원 볼륨을 표면에 투영하는 방식이라, 좁고 긴 면은 시선 각도에 따라 투영 결과가 달라짐.", body_h=1.45)
y = story_block(s, 1.5, y, 14.6, *TAG_FIX, "아티팩트 진원지 자체를 제거",
    "각도별 보정 대신 레인 채움면을 없애고, 차선 표기를 흰 구분선+황색 중앙선(픽셀 폭 선)으로 일원화 — "
    "실제 도로 표기이자 SUMO 등 시뮬레이터 표준 방식. 픽셀 폭 선은 시점과 무관하게 안정.", body_h=1.45)
story_block(s, 1.5, y, 14.6, *TAG_EFF, "형상 안정 + 메모리 절감",
    "카메라 회전 시 무늬 변형 소멸. 레인면이 지오메트리의 대부분이었어서 브라우저 메모리도 "
    "585MB → 220MB로 감소, 타일 빌드도 그만큼 빨라짐.", body_h=1.45)

card_head(s, 17.3, 3.35, 15.4, "성능 — 프레임이 한 자리로 떨어지던 원인", C_AMBER)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
y = 4.35
y = story_block(s, 17.7, y, 14.6, *TAG_CAUSE, "화면 밖 오브젝트도 CPU를 씀",
    "노드·포트·커넥션 아이콘(엔티티)은 “멀면 안 그리는” GPU 컬링을 걸어도, 존재하는 한 매 프레임 "
    "CPU 처리 대상. 팬을 하다 보면 수천 개가 누적돼 FPS가 한 자리로 하락 (4,597개 실측).", body_h=1.45)
y = story_block(s, 17.7, y, 14.6, *TAG_FIX, "거리 기반 생명주기",
    "카메라 1km 진입 시 생성, 1.3km 이탈 시 삭제(데이터는 캐시에 있어 재진입 시 재생성). "
    "“안 보이게”가 아니라 “존재 자체를” 화면 주변으로 제한해 총량을 수백 개로 유지.", body_h=1.45)
story_block(s, 17.7, y, 14.6, *TAG_EFF, "엔티티 4,597 → 645개",
    "같은 위치·같은 줌에서 엔티티 86% 감소, 힙 697→262MB. 도로 클릭 시 차선이 아닌 "
    "도로 링크가 선택되도록 픽 해석도 일원화.", body_h=1.45)

# ════════════════════════════════════════════════════════════════
# 8 — 시뮬레이션 ① 대량 생성
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "04 시뮬레이션", C_GREEN)
add_title(s, "시뮬레이션 ① — 18.6만 대를 OOM 없이 만들기",
          "차량 수 상한(1만 대)으로 회피하던 것을, 생성 파이프라인 자체를 스트리밍으로 바꿔 제한 해제")

card_head(s, 1.1, 3.35, 15.4, "생성 단계", C_RED)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
y = 4.35
y = story_block(s, 1.5, y, 14.6, *TAG_PROB, "대전·세종 전역에 차가 5대만 보임",
    "도로 연장 기반 적정 차량 수는 12,407km × 15대/km = 186,110대. 하지만 생성 시 이벤트 4,830만 건을 "
    "리스트에 전부 쌓은 뒤 DB에 쓰는 구조라 서버 힙이 터짐(OOM) → 임시로 1만 대 상한을 걸어 밀도가 1/18로.", body_h=1.45)
y = story_block(s, 1.5, y, 14.6, *TAG_FIX, "Consumer 스트리밍 생성",
    "차량 시뮬레이션 코어가 이벤트를 리스트에 모으는 대신 콜백(sink)으로 흘려보내고, sink가 SQLite에 "
    "5만 건 단위로 직접 기록. 메모리에는 항상 차량 1대분만 존재 — 이벤트 수와 무관하게 동작.", body_h=1.45)
story_block(s, 1.5, y, 14.6, *TAG_EFF, "186,110대 / 70초 / 크래시 0",
    "전 지역 정상 밀도로 생성. 진행률(“15만/18.6만대”)이 지도 상단에 실시간 표시. "
    "생성물은 대용량이므로 전체 CZML을 만들지 않고(largeMode) viewport 스트리밍이 서빙.", body_h=1.45)

card_head(s, 17.3, 3.35, 15.4, "함께 잡은 서버 이슈", C_PURPLE)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
vitems = [
    ("디스크 70GB 누수", "교통량 히트맵이 1초마다 집계 API를 부르는데, 호출마다 120MB DB 사본을 만들고 지우지 않음 — "
     "임시 사본 621개(70GB) 실측. 시나리오당 사본 1개를 캐시 재사용하도록 전면 수정"),
    ("무한 재생성 루프", "OOM은 Exception이 아닌 Error 계열이라 실패 등록을 우회 → 프론트 폴링이 생성을 무한 재시작해 "
     "OOM이 반복. Throwable 처리로 실패를 확정시켜 루프 차단"),
    ("대용량 가드", "이벤트 500만 건 초과 DB는 전체 CZML 빌드 진입 자체를 차단하고 largeMode 응답 — "
     "실수로 전체 로드를 트리거해도 서버가 죽지 않음"),
    ("의도치 않은 더미 생성", "시나리오를 열기만 해도 데이터가 없으면 자동으로 더미를 만들던 경로 분리 — "
     "존재 확인 후 로드, 생성은 명시적 버튼으로만"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 17.85, yy + 0.20, 0.13, C_PURPLE)
    txt(s, head, 18.10, yy, 14.2, 0.45, size=10.5, bold=True)
    txt(s, body, 18.10, yy + 0.44, 14.1, 1.30, size=9.5, color=C_TEXT2)
    yy += 2.05

# ════════════════════════════════════════════════════════════════
# 9 — 시뮬레이션 ② viewport+시간창 스트리밍
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "04 시뮬레이션", C_GREEN)
add_title(s, "시뮬레이션 ② — 화면과 재생 시각 기준으로만 차량을 로드",
          "차량은 시공간 데이터 — 공간(bbox)과 시간(재생 창)을 함께 잘라야 한다")

card_head(s, 1.1, 3.35, 15.4, "선별 방식 (서버)", C_TEAL)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
vitems = [
    ("“창 안에 화면을 지나는 차”", "화면 bbox 안의 도로(네트워크 RTree 재사용)를, 재생 시간창 안에 한 번이라도 "
     "지나는 차량을 선별. 화면 밖에서 진입 예정인 차도 미리 포함되어 경계에서 뿅 나타나지 않음"),
    ("궤적은 통째로", "선별된 차량은 창 안의 전체 궤적(화면 밖 구간 포함)을 반환 — 차가 화면을 나갔다 들어와도 "
     "보간(부드러운 이동)이 끊기지 않음. 창 양끝 ±60초 버퍼로 경계 보간도 연속"),
    ("붐비면 오래 보이는 차 우선", "화면당 상한(1,500대) 초과 시 화면 체류시간이 긴 차량부터 선별 + 결정적 정렬 — "
     "창이 갱신될 때 차들이 무작위로 바뀌며 깜빡이지 않음"),
    ("응답 속도", "네트워크 파싱·좌표 변환기를 시나리오당 1회 캐시 — 첫 요청 4.6초, 이후 0.9초"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 1.45, yy + 0.20, 0.13, C_TEAL)
    txt(s, head, 1.70, yy, 14.4, 0.45, size=10.5, bold=True)
    txt(s, body, 1.70, yy + 0.44, 14.3, 1.30, size=9.5, color=C_TEAL if False else C_TEXT2)
    yy += 2.05

card_head(s, 17.3, 3.35, 15.4, "재생 경험 (클라이언트)", C_GREEN)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
vitems = [
    ("끊김 없는 이어재생", "재생 시각이 창의 끝 30초 전에 도달하면 다음 창을 미리 로드. 화면을 움직여 데이터를 "
     "새로 받아도 재생 시각을 보존 — “팬하면 0초로 돌아가던” 문제 수정"),
    ("밀집 시 표현 자동 전환", "화면 안 차량이 1,500대를 넘으면 개별 차량 대신 링크별 교통량 히트맵으로 전환, "
     "줌인해서 줄어들면 개별 차량 복귀 (경계 진동 방지 히스테리시스 + 최소 4초 간격)"),
    ("항상 보이는 상태 배지", "지도 좌하단에 “차량 823대 표시 / 전체 1,204대” — 지금 화면이 전체인지 "
     "일부인지, 히트맵 모드인지 즉시 인지"),
    ("배속 회귀 수정", "재생 배속에 비례해 시간창을 무제한 확장하던 로직이 수천 대 로드→FPS 6을 유발 — "
     "창 상한 300초 + 벽시계 8초 최소 재요청 간격으로 안정화"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 17.85, yy + 0.20, 0.13, C_GREEN)
    txt(s, head, 18.10, yy, 14.2, 0.45, size=10.5, bold=True)
    txt(s, body, 18.10, yy + 0.44, 14.1, 1.30, size=9.5, color=C_TEXT2)
    yy += 2.05

# ════════════════════════════════════════════════════════════════
# 10 — 개선 과제 P1/P2
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "05 개선 과제", C_AMBER)
add_title(s, "향후 개선 과제 ① — 단기 (P1 실사용 마감 · P2 확장)",
          "P1: 이번 작업의 잔여 검증과 마감 / P2: 같은 패턴을 다른 데이터로 확장")

card_head(s, 1.1, 3.35, 15.4, "P1 — 실사용 검증·마감", C_RED)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
vitems = [
    ("3D 형상 실기기 최종 확인", "시점 의존 아티팩트 수정이 headless 검증까지 완료된 상태 — 실제 GPU/모니터 환경에서 "
     "카메라 회전·틸트 시 잔여 아티팩트 여부 최종 확인 필요"),
    ("타일 모드 편집 저장 검증", "편집 중 동결·삭제 추적·diff 저장 로직은 구현·단위검증 완료. 브라우저에서 실제 "
     "그리기→수정→저장→재로드 사이클 검증 필요 (전체 덮어쓰기 저장과의 충돌이 유일한 위험 지점)"),
    ("2D/3D 도로 표현 통일", "3D는 상·하행 분리를 위해 우측 오프셋 적용, 2D는 아직 중심선 배치 — 같은 도로가 "
     "두 지도에서 다른 위치에 보임. 2D에도 동일 오프셋 적용 여부 결정 필요"),
    ("재생 중 재로드 체감 측정", "viewport 이동 시 차량 데이터 교체(worker 전체 재빌드)의 순간 끊김 정도를 "
     "실사용 기준으로 측정 — 심하면 P3의 증분 빌드로"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 1.45, yy + 0.20, 0.13, C_RED)
    txt(s, head, 1.70, yy, 14.4, 0.45, size=10.5, bold=True)
    txt(s, body, 1.70, yy + 0.44, 14.3, 1.30, size=9.5, color=C_TEXT2)
    yy += 2.05

card_head(s, 17.3, 3.35, 15.4, "P2 — 시설물·데이터 확장", C_AMBER)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
vitems = [
    ("시설물 공용 타일 서비스", "버스/철도 정류장, 노면표시, 노선을 신호와 같은 패턴(SQLite+RTree, bbox API)으로 — "
     "신호 타일링이 이미 같은 구조로 완성되어 있어 복제 비용 낮음"),
    ("실 시뮬레이터(NEXTSIM) 연동", "지금까지는 더미 생성기로 검증 — 실제 시뮬레이터가 출력한 vehicle_sim.db "
     "스키마로 스트리밍 파이프라인 전 구간 검증"),
    ("정밀도로지도·3D Tiles", "OpenDRIVE 기반 정밀도로지도 및 건물 3D Tiles와의 표시 정합(BOTH 드레이프) 검증"),
    ("road_class 도입 대비", "현재 LOD 랭크는 차선수/제한속도 프록시 — 데이터에 도로 등급이 들어오면 "
     "프록시 로직만 교체하면 되도록 격리되어 있음"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 17.85, yy + 0.20, 0.13, C_AMBER)
    txt(s, head, 18.10, yy, 14.2, 0.45, size=10.5, bold=True)
    txt(s, body, 18.10, yy + 0.44, 14.1, 1.30, size=9.5, color=C_TEXT2)
    yy += 2.05

# ════════════════════════════════════════════════════════════════
# 11 — 개선 과제 P3 + 환경
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "05 개선 과제", C_PURPLE)
add_title(s, "향후 개선 과제 ② — 구조 개선(P3) 및 환경 이슈",
          "동작은 하지만 구조적으로 더 나은 형태가 있는 항목들 — 중장기 리팩터링 후보")

card_head(s, 1.1, 3.35, 15.4, "P3 — 구조 개선 (중장기)", C_PURPLE)
shadow_card(s, 1.1, 3.97, 15.4, 8.9)
vitems = [
    ("초기화 상태머신화", "타일/레이어 초기화의 레이스 조건들을 개별 가드(재진입 방지, 폴링, 재시도)로 잡아 둔 상태 — "
     "명시적 상태머신으로 재편하면 “조작하다 보면 엉키는 느낌”의 잔여 가능성 원천 제거"),
    ("차량 worker 증분 빌드", "viewport 이동 시 차량 보간 테이블을 전체 재빌드 — 차량 인덱스 매핑(궤적·꼬리·모델 3개 "
     "소비자)이 얽혀 있어 대형 작업이지만, 재생 중 끊김의 근본 해소책"),
    ("네트워크 3D Tiles 정적 생성", "도로를 매번 실시간 드레이프하는 대신 3D Tiles로 사전 생성 — 렌더 비용을 "
     "근본적으로 낮추고 시점 아티팩트 자유도 확보 (설계 문서 작성됨)"),
    ("LOD 경계 상수 단일화", "2D(해상도)와 3D(고도) LOD 경계 환산이 근사치로 이원화 — 단일 출처로 정리"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 1.45, yy + 0.20, 0.13, C_PURPLE)
    txt(s, head, 1.70, yy, 14.4, 0.45, size=10.5, bold=True)
    txt(s, body, 1.70, yy + 0.44, 14.3, 1.30, size=9.5, color=C_TEXT2)
    yy += 2.05

card_head(s, 17.3, 3.35, 15.4, "환경·인프라 (별도 트랙)", C_TEXT3)
shadow_card(s, 17.3, 3.97, 15.4, 8.9)
vitems = [
    ("지형 서버 인증서", "자체서명 인증서라 새 브라우저/시크릿 모드에서 지형 로드가 조용히 실패 — 3D가 검게 보이는 "
     "첫인상 문제. 신뢰 인증서 전환 또는 프록시 필요"),
    ("Cesium Ion 토큰", "만료된 토큰(401)이 콘솔 에러를 지속 발생 — 지형 폴백 경로 정리 겸 토큰 갱신"),
    ("Naver 지도 연동", "2D 배경 Naver 지도 연동 진행 중 — SDK 인증 실패가 Cesium 렌더 루프를 멈추던 문제에 "
     "가드 적용됨. 키/도메인 등록 후 활성화"),
    ("데이터 정합", "KTDB 변환 XML의 구성요소 구조가 파서 기대와 달라 포트/커넥션이 조용히 유실되던 문제 수정 — "
     "유사 불일치의 자동 검증(스키마 테스트) 추가 고려"),
]
yy = 4.35
for head, body in vitems:
    sdot(s, 17.85, yy + 0.20, 0.13, C_TEXT3)
    txt(s, head, 18.10, yy, 14.2, 0.45, size=10.5, bold=True)
    txt(s, body, 18.10, yy + 0.44, 14.1, 1.30, size=9.5, color=C_TEXT2)
    yy += 2.05

# ════════════════════════════════════════════════════════════════
# 12 — 정량 요약 + 결론
# ════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_bg(s); add_header(s, "요약", C_GREEN)
add_title(s, "정량 요약 및 결론",
          "측정 조건: 대전·세종 KTDB 실데이터 (링크 62,434 / 노드 47,181), 브라우저 E2E·API 직접 계측")

mx = [1.1, 9.3, 17.5, 25.7]
metric(s, mx[0], 3.4, 7.8, "62,434 링크", "전국급 네트워크 스트리밍\n(전체-로드 한계 ~10만 링크 돌파)", C_BLUE)
metric(s, mx[1], 3.4, 7.8, "186,110 대", "차량 스트리밍 생성 70초\n(기존: 서버 OOM 크래시)", C_TEAL)
metric(s, mx[2], 3.4, 7.8, "0.9 초", "viewport 차량 응답(캐시 후)\n첫 요청 4.6초", C_GREEN)
metric(s, mx[3], 3.4, 7.8, "220 MB", "JS Heap (기존 ~1GB+)\n엔티티 4,597 → 645개", C_PURPLE)

metric(s, mx[0], 5.95, 7.8, "~17 초", "시나리오 진입→지도 표시\n(106MB 일괄 다운로드 제거)", C_BLUE)
metric(s, mx[1], 5.95, 7.8, "70 GB", "서버 디스크 누수 회수\n(DB 사본 캐시화)", C_RED)
metric(s, mx[2], 5.95, 7.8, "5 단계", "LOD (macro~detail)\n+ 밀집 히트맵 자동 전환", C_AMBER)
metric(s, mx[3], 5.95, 7.8, "100 %", "additive API 설계\n기존 경로 무중단 병존", C_TEAL)

card_head(s, 1.1, 8.75, 32.0, "결론", C_GREEN)
shadow_card(s, 1.1, 9.37, 32.0, 3.35)
txt(s, "“전부 내려받아 전부 그린다” 구조를 “보이는 것만, 보이는 수준으로” 스트리밍하는 구조로 전환했다. "
        "공간(타일)·디테일(LOD 5단계)·시간(재생 창) 세 축으로 데이터를 자르고, 서버 공간 인덱스가 그 계산을 담당한다.",
    1.5, 9.70, 31.2, 1.1, size=11, color=C_TEXT)
txt(s, "이로써 단일 도시(~10만 링크)가 한계였던 플랫폼이 광역권 전체(대전·세종, 링크 6.2만 + 차량 18.6만 대)를 "
        "단일 브라우저 세션에서 구동함을 실증했다. 모든 신규 API는 기존 경로와 병존(additive)하므로 무중단 적용이 가능하며, "
        "동일 패턴을 신호·시설물·실 시뮬레이터 결과로 확장하는 기반이 마련되었다.",
    1.5, 10.85, 31.2, 1.7, size=11, color=C_TEXT)

out = "/Users/hskim/Documents/repo/iitp/IITP_네트워크타일링_변경내역_개선과제.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
