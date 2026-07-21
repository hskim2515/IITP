"""
IITP 2026 연구사업계획 PPT
참고 디자인: 가이아쓰리디 스타일
- 흰 배경
- 상단 회색(ADB3C1)+골드(FDC951) 헤더바
- 골드(F9BE36/FFC000) 타이틀 배너
- 다크슬레이트(383E4C) 강조박스
- 레드(FF0000) 포인트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from lxml import etree
import copy, math

# ── 색상 ─────────────────────────────────────────────────────
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BG      = RGBColor(0xFF, 0xFF, 0xFF)   # 흰 배경
C_HEADER1 = RGBColor(0xAD, 0xB3, 0xC1)  # 헤더 회색
C_HEADER2 = RGBColor(0xFD, 0xC9, 0x51)  # 헤더 골드
C_GOLD1   = RGBColor(0xF9, 0xBE, 0x36)  # 타이틀 배너 골드
C_GOLD2   = RGBColor(0xFF, 0xC0, 0x00)  # 타이틀 배너 진골드
C_DARK    = RGBColor(0x38, 0x3E, 0x4C)  # 다크 슬레이트
C_RED     = RGBColor(0xFF, 0x00, 0x00)  # 레드 포인트
C_BLACK   = RGBColor(0x00, 0x00, 0x00)
C_GRAY    = RGBColor(0x66, 0x6E, 0x7A)
C_LGRAY   = RGBColor(0xF0, 0xF2, 0xF5)  # 밝은 회색 배경
C_LGOLD   = RGBColor(0xFF, 0xF0, 0xB0)  # 연한 골드 배경
C_DRED    = RGBColor(0xCC, 0x00, 0x00)  # 진한 레드
C_BLUE    = RGBColor(0x1E, 0x6F, 0xBB)  # 파란 보조색
C_GREEN   = RGBColor(0x2E, 0x9E, 0x6B)  # 녹색 보조색

prs = Presentation()
prs.slide_width  = Cm(33.87)
prs.slide_height = Cm(19.05)
blank = prs.slide_layouts[6]

W = 33.87   # cm
H = 19.05

# ═══════════════════════════════════════════════════════════════
# 헬퍼
# ═══════════════════════════════════════════════════════════════
def cm(v): return Cm(v)

def rect(slide, l, t, w, h, fill, line=None, lw=0.5, r=0):
    shape = slide.shapes.add_shape(
        1 if r == 0 else 5,
        cm(l), cm(t), cm(w), cm(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(lw)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=12, bold=False,
        color=C_BLACK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(cm(l), cm(t), cm(w), cm(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def txt_ml(slide, lines, l, t, w, h, default_size=11,
           default_color=C_BLACK, default_bold=False, wrap=True):
    """lines: list of str or (text, size, bold, color)"""
    txb = slide.shapes.add_textbox(cm(l), cm(t), cm(w), cm(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        if isinstance(item, str):
            s, sz, bd, col = item, default_size, default_bold, default_color
        else:
            s   = item[0]
            sz  = item[1] if len(item) > 1 else default_size
            bd  = item[2] if len(item) > 2 else default_bold
            col = item[3] if len(item) > 3 else default_color
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = s
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = col
    return txb

def circle(slide, l, t, d, fill, line=None):
    """원형 도형"""
    shape = slide.shapes.add_shape(9, cm(l), cm(t), cm(d), cm(d))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line: shape.line.color.rgb = line
    else: shape.line.fill.background()
    return shape

def add_header(slide, company="가이아쓰리디"):
    """공통 상단 헤더 바 (회색 + 골드)"""
    # 회색 전체 바
    r = rect(slide, 0, 0, W, 2.3, C_HEADER1)
    # 골드 우측 1/3
    rect(slide, W*0.65, 0, W*0.35, 2.3, C_HEADER2)
    # 흰 대각선 분리선 효과 (얇은 사다리꼴처럼 보이게)
    rect(slide, W*0.62, 0, 0.35, 2.3, C_WHITE)
    rect(slide, W*0.635, 0, 0.2, 2.3, C_HEADER1)
    # 회사명
    txt(slide, f"[{company}]", 1.2, 0.75, 6.0, 0.85,
        size=14, bold=False, color=C_WHITE)
    return r

def add_title_banner(slide, title_main, title_sub=""):
    """골드 타이틀 배너"""
    rect(slide, 0.8, 2.5, W-1.6, 1.9, C_GOLD1)
    rect(slide, 0.8, 2.5, 0.45, 1.9, C_GOLD2)
    txt(slide, title_main, 1.55, 2.58, W-3.5, 1.0,
        size=22, bold=True, color=C_DARK)
    if title_sub:
        txt(slide, title_sub, 1.55, 3.55, W-3.5, 0.7,
            size=12, color=C_DARK)

def numbered_box(slide, num, l, t, d=1.1, bg=C_RED, fg=C_WHITE):
    """번호 원형 박스"""
    circle(slide, l, t, d, bg)
    txt(slide, str(num), l+0.05, t+0.1, d, d-0.2,
        size=20, bold=True, color=fg, align=PP_ALIGN.CENTER)

def flow_arrow(slide, l, t, size=0.5):
    """→ 화살표 도형"""
    shape = slide.shapes.add_shape(13, cm(l), cm(t), cm(size*1.5), cm(size))
    shape.fill.solid(); shape.fill.fore_color.rgb = C_RED
    shape.line.fill.background()

def stat_box(slide, l, t, w, h, num, unit, label, bg=C_DARK, fg=C_GOLD1):
    """KPI 수치 박스 인포그래픽"""
    rect(slide, l, t, w, h, bg)
    txt(slide, num,   l+0.2, t+0.25, w-0.4, h*0.45,
        size=32, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txt(slide, unit,  l+0.2, t+h*0.52, w-0.4, h*0.2,
        size=11, color=C_LGOLD, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.2, t+h*0.73, w-0.4, h*0.25,
        size=9.5, color=C_WHITE, align=PP_ALIGN.CENTER)

def card_with_top(slide, l, t, w, h, top_color, top_h, title, items, dot_color=None):
    """상단 컬러 + 내용 카드"""
    if dot_color is None: dot_color = top_color
    rect(slide, l, t, w, h, C_LGRAY)
    rect(slide, l, t, w, top_h, top_color)
    txt(slide, title, l+0.3, t+0.08, w-0.6, top_h-0.12,
        size=12, bold=True, color=C_WHITE)
    for i, item in enumerate(items):
        # 불릿 사각형
        r2 = slide.shapes.add_shape(1, cm(l+0.3), cm(t+top_h+0.3+i*0.72),
                                     cm(0.18), cm(0.18))
        r2.fill.solid(); r2.fill.fore_color.rgb = dot_color
        r2.line.fill.background()
        txt(slide, item, l+0.62, t+top_h+0.22+i*0.72,
            w-0.8, 0.65, size=10.5, color=C_BLACK)

def process_step(slide, steps, l, t, step_w, step_h, colors=None):
    """가로 프로세스 흐름도"""
    if colors is None:
        colors = [C_DARK, C_GOLD1, C_RED, C_BLUE, C_GREEN]
    gap = 0.5
    total_w = len(steps) * step_w + (len(steps)-1) * gap
    for i, (title, desc) in enumerate(steps):
        lx = l + i*(step_w + gap)
        c = colors[i % len(colors)]
        rect(slide, lx, t, step_w, step_h, c)
        # 번호
        circle(slide, lx+0.15, t+0.15, 0.6, C_WHITE)
        txt(slide, str(i+1), lx+0.15, t+0.2, 0.62, 0.5,
            size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        txt(slide, title, lx+0.9, t+0.18, step_w-1.1, 0.7,
            size=11, bold=True, color=C_WHITE)
        # 설명
        for di, dl in enumerate(desc.split('\n')):
            txt(slide, dl, lx+0.2, t+0.95+di*0.52,
                step_w-0.4, 0.5, size=9.5, color=C_WHITE)
        # 화살표 (마지막 제외)
        if i < len(steps)-1:
            arx = lx + step_w + 0.08
            ary = t + step_h/2 - 0.22
            shape = slide.shapes.add_shape(13, cm(arx), cm(ary), cm(0.36), cm(0.44))
            shape.fill.solid(); shape.fill.fore_color.rgb = C_GRAY
            shape.line.fill.background()

# ═══════════════════════════════════════════════════════════════
# 슬라이드 01 ── 표지
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)

# 배경
rect(sl, 0, 0, W, H, C_WHITE)

# 상단 헤더
add_header(sl)

# 중앙 강조 배너 (골드)
rect(sl, 0, 5.5, W, 7.5, C_GOLD1)
rect(sl, 0, 5.5, 0.7, 7.5, C_GOLD2)

# 메인 타이틀
txt(sl, "도시교통 디지털트윈", 2.0, 5.9, 28, 2.4,
    size=46, bold=True, color=C_DARK)
txt(sl, "시뮬레이션 플랫폼 고도화 연구", 2.0, 8.0, 28, 1.8,
    size=36, bold=True, color=C_WHITE)

# 부제
rect(sl, 1.5, 10.0, 18.0, 0.8, C_DARK)
txt(sl, "네트워크 편집  ·  실시간 시뮬레이션  ·  지능형 분석  ·  운영 자동화",
    1.7, 10.05, 17.5, 0.72,
    size=13, color=C_GOLD1, align=PP_ALIGN.LEFT)

# 우측 인포그래픽: 핵심 수치
kpis = [("4개", "세부과제"), ("12+", "핵심기능"), ("2026", "연구연도")]
for i, (n, l2) in enumerate(kpis):
    bx = 22.5 + i*3.7
    stat_box(sl, bx, 6.2, 3.3, 2.8, n, "", l2, bg=C_DARK)

# 하단 정보
rect(sl, 0, 13.3, W, 5.75, C_LGRAY)
rect(sl, 0, 13.3, W, 0.06, C_GOLD2)
txt(sl, "한국도로공사 도로교통연구원   |   IITP 지원과제",
    2.0, 14.0, 25, 0.7, size=13, color=C_GRAY)
txt(sl, "2026. 03", 2.0, 14.8, 10, 0.6, size=12, color=C_GRAY)

# 하단 골드 포인트 바
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 02 ── 목차
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
add_title_banner(sl, "Contents", "연구 추진 체계 및 세부 과제")

chapters = [
    ("01", "연구 배경 및 필요성",       "교통 혼잡·디지털트윈 정책·기존 한계"),
    ("02", "기존 플랫폼 현황",          "핵심 구현 기능 8가지 요약"),
    ("03", "세부과제 1  네트워크 편집 고도화", "위상검증·OSM임포트·Diff시각화"),
    ("04", "세부과제 2  시뮬레이션 고도화",   "실시간연동·이벤트편집기·A/B비교"),
    ("05", "세부과제 3  분석 기능 확장",  "교차로성능·속도프로파일·탄소배출"),
    ("06", "세부과제 4  UX 및 운영 편의", "영상내보내기·자동리포트·협업기능"),
    ("07", "추진 일정 (로드맵)",         "분기별 산출물 및 마일스톤"),
]

# 목차 아이템: 2열 배치
for i, (num, title, desc) in enumerate(chapters):
    col = i // 4
    row = i %  4
    lx = 0.9 + col*16.8
    ty = 4.65 + row * 3.45

    # 박스
    rect(sl, lx, ty, 16.0, 3.0, C_LGRAY, line=C_HEADER1, lw=0.5)
    # 좌측 번호 세로 바
    rect(sl, lx, ty, 1.4, 3.0, C_DARK)
    txt(sl, num, lx+0.05, ty+0.6, 1.35, 1.5,
        size=20, bold=True, color=C_GOLD1, align=PP_ALIGN.CENTER)

    # 내용
    txt(sl, title, lx+1.65, ty+0.35, 14.0, 0.9,
        size=13.5, bold=True, color=C_DARK)
    rect(sl, lx+1.65, ty+1.3, 12.5, 0.05, C_GOLD2)
    txt(sl, desc, lx+1.65, ty+1.5, 13.8, 0.8,
        size=10.5, color=C_GRAY)

# 우측 빈 공간 장식
rect(sl, 0.9+16.8, 4.65, 1.0, H-5.0, C_LGRAY)  # 배경만

# ═══════════════════════════════════════════════════════════════
# 슬라이드 03 ── 연구 배경 및 필요성
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
add_title_banner(sl, "연구 배경 및 필요성",
                 "01  현황과 한계, 추진 근거")

# ── 상단: 3대 배경 카드 ──
cards3 = [
    (C_DARK,  "교통 혼잡 사회적 비용",
     ["연간 약 73조 원 (2023년 기준)",
      "수도권 집중 현상 지속 심화",
      "C-ITS·자율주행 확산으로",
      "정밀 교통 분석 수요 급증"]),
    (C_RED,   "디지털트윈 정책 환경",
     ["국토부 스마트시티 디지털트윈",
      "교통 분야 국가 R&D 본격 추진",
      "공간정보 오픈 API 확대 정책",
      "실시간 교통 연계 인프라 구축"]),
    (C_BLUE,  "기존 플랫폼 한계",
     ["정적 CZML 재생 방식만 지원",
      "편집 후 위상 오류 자동감지 없음",
      "외부 데이터(OSM/SHP) 임포트 불가",
      "분석 결과 보고서 자동화 미지원"]),
]
for j, (col, title, items) in enumerate(cards3):
    lx = 0.9 + j * 10.9
    card_with_top(sl, lx, 4.7, 10.5, 5.8, col, 0.7, title, items, col)

# ── 하단: 연구 방향 프로세스 ──
txt(sl, "▶  본 연구의 추진 방향", 0.9, 10.85, 15, 0.6,
    size=13, bold=True, color=C_DARK)
rect(sl, 0.9, 11.55, W-1.8, 0.05, C_GOLD2)

steps = [
    ("현황 분석",  "플랫폼 기능\n한계 진단"),
    ("과제 도출",  "4대 세부\n과제 정의"),
    ("설계·구현",  "기능별\n개발 수행"),
    ("검증·테스트","위상·성능\n품질 확인"),
    ("성과 배포",  "보고서·운영\n체계 완성"),
]
process_step(sl, steps, l=0.9, t=12.0, step_w=6.1, step_h=2.5,
             colors=[C_DARK, C_GOLD1, C_RED, C_BLUE, C_GREEN])

# 하단 골드 바
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 04 ── 기존 플랫폼 현황
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
add_title_banner(sl, "기존 플랫폼 핵심 기능 현황",
                 "02  현재 구현 완료된 핵심 모듈")

features = [
    (C_DARK,  "도로 네트워크 편집",
     "Node · Link · Lane · Connection\n스키마 기반 동적 UI · 히스토리 관리"),
    (C_RED,   "차량 시뮬레이션",
     "CZML + Web Worker 비동기 처리\n2D(OL) + 3D(Cesium) 동시 뷰"),
    (C_BLUE,  "교통 제어 편집",
     "신호 타이밍 · 노면마킹\n버스·철도 정류장 관리"),
    (C_GREEN, "분석 레이어",
     "히트맵 · OD 매트릭스\nTrail 분석 · Recharts 통계"),
]

# 2×2 그리드 카드
positions = [(0.9,4.6),(17.45,4.6),(0.9,11.4),(17.45,11.4)]
for (lx, ty), (col, title, desc) in zip(positions, features):
    w_card = 16.0
    h_card = 6.3
    rect(sl, lx, ty, w_card, h_card, C_LGRAY, line=col, lw=1.0)
    rect(sl, lx, ty, w_card, 1.0, col)

    # 번호 원
    idx = positions.index((lx,ty))+1
    numbered_box(sl, f"0{idx}", lx+0.25, ty+0.15, d=0.72, bg=C_WHITE, fg=col)

    txt(sl, title, lx+1.2, ty+0.15, w_card-1.5, 0.75,
        size=14, bold=True, color=C_WHITE)

    for di, line in enumerate(desc.split('\n')):
        rect(sl, lx+0.35, ty+1.35+di*0.95, 0.22, 0.22, col)
        txt(sl, line, lx+0.7, ty+1.28+di*0.95,
            w_card-1.1, 0.75, size=12.5, color=C_BLACK)

# 하단 기술 스택
rect(sl, 0, H-1.3, W, 0.95, C_DARK)
txt(sl, "Tech Stack :  Spring Boot  ·  PostgreSQL  ·  React / TypeScript  ·  "
        "CesiumJS  ·  OpenLayers  ·  Web Workers  ·  Recharts",
    1.0, H-1.22, W-2.0, 0.75,
    size=11, color=C_GOLD1, align=PP_ALIGN.LEFT)
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 05 ── 세부과제 1: 네트워크 편집 고도화
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)

# 세부과제 배지
rect(sl, 0.9, 2.55, 3.5, 0.8, C_DARK)
txt(sl, "세부과제 01", 0.95, 2.65, 3.4, 0.65,
    size=11, bold=True, color=C_GOLD1, align=PP_ALIGN.CENTER)

add_title_banner(sl, "네트워크 편집 고도화",
                 "03  위상 검증 · OSM 임포트 · 버전 Diff 시각화")

# 3개 세부 기능 카드
sub_tasks = [
    (C_DARK,  "① 위상 검증 (Topology Validation)",
     ["단절 링크·고립 노드·차선 방향 충돌 자동 감지",
      "중복 좌표·역방향 링크 실시간 검출 엔진",
      "오류 요소 지도 하이라이트 + 원클릭 이동",
      "저장 전 검증 게이트 → 데이터 품질 보증",
      "오류 유형별 통계 패널 제공"]),
    (C_RED,   "② OSM / SHP 외부 데이터 임포트",
     ["OpenStreetMap Overpass API 실시간 연동",
      "SHP(Shapefile) 파일 업로드 및 변환 지원",
      "좌표계 자동 변환 (WGS84 ↔ EPSG:5186)",
      "임포트 후 자동 위상 정제 파이프라인",
      "현장 조사 없이 초기 네트워크 신속 구축"]),
    (C_BLUE,  "③ 버전 간 네트워크 Diff 시각화",
     ["v1 → v2 변경 사항 색상 구분 (추가/삭제/수정)",
      "링크·노드·차선별 변경 이력 패널 표시",
      "세션 이력 → 서버 영구 저장 이력으로 확장",
      "특정 버전 원클릭 복원 기능",
      "변경 요약 자동 리포트 생성"]),
]
for j, (col, title, items) in enumerate(sub_tasks):
    lx = 0.9 + j * 10.9
    ty = 4.65
    h_card = 8.5
    rect(sl, lx, ty, 10.5, h_card, C_LGRAY, line=col, lw=0.8)
    rect(sl, lx, ty, 10.5, 0.75, col)
    txt(sl, title, lx+0.3, ty+0.1, 9.9, 0.6,
        size=12, bold=True, color=C_WHITE)
    for i, item in enumerate(items):
        # 번호 박스
        r2 = sl.shapes.add_shape(1, cm(lx+0.25), cm(ty+1.05+i*1.42), cm(0.5), cm(0.5))
        r2.fill.solid(); r2.fill.fore_color.rgb = col
        r2.line.fill.background()
        t2 = sl.shapes.add_textbox(cm(lx+0.27), cm(ty+1.07+i*1.42), cm(0.5), cm(0.45))
        t2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = t2.text_frame.paragraphs[0].add_run()
        run.text = str(i+1)
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = C_WHITE
        txt(sl, item, lx+0.88, ty+1.05+i*1.42,
            9.3, 1.25, size=11, color=C_BLACK)

# 하단 기대효과 바
rect(sl, 0, 13.5, W, 1.5, C_DARK)
txt(sl, "기대 효과", 1.0, 13.6, 4.5, 0.5,
    size=13, bold=True, color=C_GOLD1)
effects = ["데이터 품질 40% 향상", "초기 구축 시간 60% 단축", "변경 이력 완전 추적"]
for i, e in enumerate(effects):
    rect(sl, 5.5+i*9.3, 13.62, 8.8, 1.2, C_GOLD1)
    txt(sl, e, 5.7+i*9.3, 13.75, 8.4, 0.8,
        size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 06 ── 세부과제 2: 시뮬레이션 고도화
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
rect(sl, 0.9, 2.55, 3.5, 0.8, C_RED)
txt(sl, "세부과제 02", 0.95, 2.65, 3.4, 0.65,
    size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_title_banner(sl, "시뮬레이션 고도화",
                 "04  실시간 연동 · 이벤트 편집기 · A/B 비교")

sub_sim = [
    (C_DARK,  "① 실시간 시뮬레이션 연동",
     ["WebSocket / SSE 기반 스트리밍 수신",
      "SUMO·VISSIM 시뮬레이터 서버 연동",
      "진행 중 시뮬레이션 라이브 모니터링",
      "재생·정지·배속 실시간 제어 UI",
      "연결 끊김 자동 재연결 및 버퍼 관리"]),
    (C_RED,   "② 시나리오 이벤트 편집기",
     ["교통사고·공사·신호고장 타임라인 삽입",
      "드래그 앤 드롭 이벤트 배치 UI",
      "이벤트 발생 전/후 교통 흐름 비교",
      "이벤트 연쇄 효과 (도미노 모델) 지원",
      "시나리오 템플릿 저장·공유 기능"]),
    (C_BLUE,  "③ A/B 시나리오 비교 뷰",
     ["화면 분할(Split-view) 동기 재생",
      "KPI 나란히 표시 (지체·속도·대기행렬)",
      "시나리오 간 차이 히트맵 오버레이",
      "최적 시나리오 자동 추천 알고리즘",
      "비교 결과 PDF/CSV 내보내기"]),
]
for j, (col, title, items) in enumerate(sub_sim):
    lx = 0.9 + j * 10.9
    ty = 4.65
    h_card = 8.5
    rect(sl, lx, ty, 10.5, h_card, C_LGRAY, line=col, lw=0.8)
    rect(sl, lx, ty, 10.5, 0.75, col)
    txt(sl, title, lx+0.3, ty+0.1, 9.9, 0.6,
        size=12, bold=True, color=C_WHITE)
    for i, item in enumerate(items):
        r2 = sl.shapes.add_shape(1, cm(lx+0.25), cm(ty+1.05+i*1.42), cm(0.5), cm(0.5))
        r2.fill.solid(); r2.fill.fore_color.rgb = col
        r2.line.fill.background()
        t2 = sl.shapes.add_textbox(cm(lx+0.27), cm(ty+1.07+i*1.42), cm(0.5), cm(0.45))
        t2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = t2.text_frame.paragraphs[0].add_run()
        run.text = str(i+1)
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = C_WHITE
        txt(sl, item, lx+0.88, ty+1.05+i*1.42, 9.3, 1.25, size=11, color=C_BLACK)

rect(sl, 0, 13.5, W, 1.5, C_DARK)
txt(sl, "기대 효과", 1.0, 13.6, 4.5, 0.5, size=13, bold=True, color=C_GOLD1)
effects = ["라이브 모니터링 실현", "정책 시나리오 검증 효율화", "의사결정 지원 고도화"]
for i, e in enumerate(effects):
    rect(sl, 5.5+i*9.3, 13.62, 8.8, 1.2, C_GOLD1)
    txt(sl, e, 5.7+i*9.3, 13.75, 8.4, 0.8,
        size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 07 ── 세부과제 3: 분석 기능 확장
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
rect(sl, 0.9, 2.55, 3.5, 0.8, C_BLUE)
txt(sl, "세부과제 03", 0.95, 2.65, 3.4, 0.65,
    size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_title_banner(sl, "분석 기능 확장",
                 "05  교차로 성능 · 속도 프로파일 · 탄소 배출")

sub_ana = [
    (C_DARK,  "① 교차로 성능 분석",
     ["포화도(v/c ratio) 자동 산출 엔진",
      "평균 지체 및 최대 대기행렬 계산",
      "신호 현시 최적화 제안 알고리즘",
      "HCM 기준 서비스 수준(LOS) 등급화",
      "교차로 유형별 비교 벤치마크"]),
    (C_RED,   "② 구간 통행시간 / 속도 프로파일",
     ["경로 선택 → 시간대별 속도 그래프",
      "Recharts 인터랙티브 시각화 (기존 활용)",
      "Peak / Off-peak 자동 구분 분석",
      "링크 속도 이상치 감지 알림",
      "FCD 실측 데이터와 시뮬레이션 비교"]),
    (C_BLUE,  "③ 탄소 배출 추정 모듈",
     ["차량 속도 프로파일 → 배출계수 적용",
      "MOVES / COPERT 배출 모델 내장",
      "차종별 CO₂·NOx·PM2.5 산출",
      "정책 시나리오별 저감 효과 정량 비교",
      "교통영향평가 보고서 연동"]),
]
for j, (col, title, items) in enumerate(sub_ana):
    lx = 0.9 + j * 10.9
    ty = 4.65
    h_card = 8.5
    rect(sl, lx, ty, 10.5, h_card, C_LGRAY, line=col, lw=0.8)
    rect(sl, lx, ty, 10.5, 0.75, col)
    txt(sl, title, lx+0.3, ty+0.1, 9.9, 0.6,
        size=12, bold=True, color=C_WHITE)
    for i, item in enumerate(items):
        r2 = sl.shapes.add_shape(1, cm(lx+0.25), cm(ty+1.05+i*1.42), cm(0.5), cm(0.5))
        r2.fill.solid(); r2.fill.fore_color.rgb = col
        r2.line.fill.background()
        t2 = sl.shapes.add_textbox(cm(lx+0.27), cm(ty+1.07+i*1.42), cm(0.5), cm(0.45))
        t2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = t2.text_frame.paragraphs[0].add_run()
        run.text = str(i+1)
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = C_WHITE
        txt(sl, item, lx+0.88, ty+1.05+i*1.42, 9.3, 1.25, size=11, color=C_BLACK)

rect(sl, 0, 13.5, W, 1.5, C_DARK)
txt(sl, "기대 효과", 1.0, 13.6, 4.5, 0.5, size=13, bold=True, color=C_GOLD1)
effects = ["교차로 운영 효율 15% 향상", "탄소 저감 정량화", "HCM 기준 LOS 자동 검증"]
for i, e in enumerate(effects):
    rect(sl, 5.5+i*9.3, 13.62, 8.8, 1.2, C_GOLD1)
    txt(sl, e, 5.7+i*9.3, 13.75, 8.4, 0.8,
        size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 08 ── 세부과제 4: UX 및 운영
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
rect(sl, 0.9, 2.55, 3.5, 0.8, C_GREEN)
txt(sl, "세부과제 04", 0.95, 2.65, 3.4, 0.65,
    size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_title_banner(sl, "UX 및 운영 편의 향상",
                 "06  영상 내보내기 · 자동 리포트 · 협업 기능")

sub_ux = [
    (C_DARK,  "① 시뮬레이션 영상 내보내기",
     ["Cesium scene.postRender 프레임 캡처",
      "ffmpeg.wasm 브라우저 내 MP4 인코딩",
      "구간 선택 (시작~종료 타임코드) 내보내기",
      "해상도·프레임레이트 사용자 설정",
      "발표·보고서·소셜미디어 공유 지원"]),
    (C_RED,   "② 대시보드 & 자동 리포트",
     ["시뮬레이션 종료 후 KPI 요약 화면",
      "PDF 자동 생성 (교통영향평가 형식)",
      "차트·지도·수치 통합 레이아웃",
      "정기 배치 리포트 스케줄링",
      "이메일·슬랙 자동 발송 연동"]),
    (C_BLUE,  "③ 협업 및 접근성 강화",
     ["다중 사용자 동시 편집 (CRDT 기반)",
      "역할별 권한 관리 (뷰어/편집/관리자)",
      "변경 사항 실시간 댓글·리뷰 기능",
      "반응형 UI (태블릿 현장 조사 지원)",
      "WCAG 2.1 접근성 기준 준수"]),
]
for j, (col, title, items) in enumerate(sub_ux):
    lx = 0.9 + j * 10.9
    ty = 4.65
    h_card = 8.5
    rect(sl, lx, ty, 10.5, h_card, C_LGRAY, line=col, lw=0.8)
    rect(sl, lx, ty, 10.5, 0.75, col)
    txt(sl, title, lx+0.3, ty+0.1, 9.9, 0.6,
        size=12, bold=True, color=C_WHITE)
    for i, item in enumerate(items):
        r2 = sl.shapes.add_shape(1, cm(lx+0.25), cm(ty+1.05+i*1.42), cm(0.5), cm(0.5))
        r2.fill.solid(); r2.fill.fore_color.rgb = col
        r2.line.fill.background()
        t2 = sl.shapes.add_textbox(cm(lx+0.27), cm(ty+1.07+i*1.42), cm(0.5), cm(0.45))
        t2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = t2.text_frame.paragraphs[0].add_run()
        run.text = str(i+1)
        run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = C_WHITE
        txt(sl, item, lx+0.88, ty+1.05+i*1.42, 9.3, 1.25, size=11, color=C_BLACK)

rect(sl, 0, 13.5, W, 1.5, C_DARK)
txt(sl, "기대 효과", 1.0, 13.6, 4.5, 0.5, size=13, bold=True, color=C_GOLD1)
effects = ["보고서 작성 시간 70% 단축", "현장 활용성 극대화", "기관 간 협업 체계 구축"]
for i, e in enumerate(effects):
    rect(sl, 5.5+i*9.3, 13.62, 8.8, 1.2, C_GOLD1)
    txt(sl, e, 5.7+i*9.3, 13.75, 8.4, 0.8,
        size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 09 ── 추진 일정 인포그래픽 로드맵
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
add_title_banner(sl, "연구 추진 일정", "07  분기별 로드맵 및 주요 마일스톤")

# 분기 헤더 바
quarters = [("Q1\n1~3월","기반 설계"), ("Q2\n4~6월","핵심 개발"),
            ("Q3\n7~9월","기능 완성"), ("Q4\n10~12월","검증·배포")]
q_colors = [C_DARK, C_RED, C_BLUE, C_GREEN]
col_w = 7.5
gap   = 0.25
for qi, ((q, sub), qc) in enumerate(zip(quarters, q_colors)):
    lx = 0.9 + qi*(col_w + gap)
    # 분기 헤더
    rect(sl, lx, 4.65, col_w, 1.1, qc)
    txt(sl, q.split('\n')[0], lx+0.2, 4.7, col_w-0.4, 0.65,
        size=18, bold=True, color=C_WHITE)
    txt(sl, q.split('\n')[1], lx+2.8, 4.72, col_w-3.0, 0.45,
        size=10, color=C_WHITE)
    txt(sl, sub, lx+0.2, 5.25, col_w-0.4, 0.4,
        size=10.5, color=C_WHITE)

# 과제 행
roadmap = [
    ("네트워크 편집", C_DARK, [
        "위상검증 설계\n고립노드 감지 구현",
        "OSM 임포트\n좌표계 변환 개발",
        "Diff 시각화\n버전 복원 기능",
        "통합 테스트\n품질 검증 완료"]),
    ("시뮬레이션", C_RED, [
        "WebSocket 설계\n아키텍처 정의",
        "실시간 연동\n이벤트 편집기",
        "A/B 비교뷰\n분할화면 개발",
        "시뮬레이터 연동\n안정화 완료"]),
    ("분석 기능", C_BLUE, [
        "교차로 분석\n설계·프로토타입",
        "속도 프로파일\nFCD 연동 개발",
        "탄소 배출 모듈\nMOVES 적용",
        "HCM 검증\n성과 측정"]),
    ("UX·리포트", C_GREEN, [
        "영상 캡처\n기반 개발",
        "대시보드 UI\nKPI 시각화",
        "PDF 자동화\n리포트 엔진",
        "배포·운영\n사용자 테스트"]),
]

for ri, (name, rc, tasks) in enumerate(roadmap):
    ty = 6.2 + ri * 2.75
    # 과제명 셀
    rect(sl, 0.0, ty, 0.9, 2.5, rc)
    t2 = sl.shapes.add_textbox(cm(0.05), cm(ty+0.3), cm(0.8), cm(1.8))
    tf = t2.text_frame; tf.word_wrap = True
    for ch in name:
        p = tf.add_paragraph() if ch != name[0] else tf.paragraphs[0]
        run = p.add_run()
        run.text = ch
        run.font.size = Pt(8.5)
        run.font.bold = True
        run.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

    for qi, task in enumerate(tasks):
        lx = 0.9 + qi*(col_w + gap)
        rect(sl, lx+0.08, ty+0.12, col_w-0.16, 2.25, C_LGRAY, line=rc, lw=0.5)
        # 좌측 컬러 세로선
        rect(sl, lx+0.08, ty+0.12, 0.2, 2.25, rc)
        for li, line in enumerate(task.split('\n')):
            txt(sl, line, lx+0.45, ty+0.35+li*0.85,
                col_w-0.7, 0.8, size=10.5, color=C_BLACK, bold=(li==0))

# 마일스톤 표시
rect(sl, 0.9+(col_w+gap)*1+col_w*0.5-0.15, 4.5, 0.3, 17.0-4.5-0.35, C_RED)
txt(sl, "▼ 중간 점검", 0.9+(col_w+gap)*1+col_w*0.5-0.8, 17.1, 3.0, 0.5,
    size=9, color=C_RED, bold=True)

rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 슬라이드 10 ── 기대 효과 & 결론 인포그래픽
# ═══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
rect(sl, 0, 0, W, H, C_WHITE)
add_header(sl)
add_title_banner(sl, "기대 효과 및 활용 방안",
                 "연구 성과 및 사회·경제적 파급효과")

# ── 상단: KPI 수치 인포그래픽 ──
kpi_data = [
    ("40%", "네트워크\n데이터 품질 향상", C_DARK),
    ("60%", "초기 구축\n시간 단축",       C_RED),
    ("70%", "보고서 작성\n비용 절감",     C_BLUE),
    ("15%", "교차로 운영\n효율 향상",     C_GREEN),
    ("<500ms", "실시간\n응답 지연",       C_GOLD1),
]
for i, (num, label, col) in enumerate(kpi_data):
    lx = 0.9 + i * 6.4
    stat_box(sl, lx, 4.7, 6.0, 3.2, num, "", label, bg=col,
             fg=C_WHITE if col != C_GOLD1 else C_DARK)

# ── 중단: 3열 활용 방안 ──
txt(sl, "▶  주요 활용 방안", 0.9, 8.3, 15, 0.6,
    size=13, bold=True, color=C_DARK)
rect(sl, 0.9, 9.0, W-1.8, 0.05, C_GOLD2)

uses = [
    (C_DARK,  "행정·정책 활용",
     ["교통영향평가 보고서 자동화",
      "지자체·공공기관 교통운영 의사결정",
      "도로 신설·개선 사업 사전 검토"]),
    (C_RED,   "산업·연구 활용",
     ["C-ITS·자율주행 시뮬레이션 환경",
      "대학·연구기관 교육·연구 플랫폼",
      "민간 교통 컨설팅 기관 도구 제공"]),
    (C_BLUE,  "환경·지속가능성",
     ["탄소 저감 시나리오 정량 비교",
      "친환경 교통정책 효과 입증 근거",
      "국가 온실가스 감축 목표 기여"]),
]
for j, (col, title, items) in enumerate(uses):
    lx = 0.9 + j * 10.9
    card_with_top(sl, lx, 9.3, 10.5, 4.8, col, 0.65, title, items, col)

# ── 하단 강조 메시지 ──
rect(sl, 0, H-1.3, W, 0.95, C_DARK)
txt(sl, "본 연구를 통해 대한민국 교통 디지털트윈의 표준 플랫폼을 구축합니다",
    1.0, H-1.22, W-2.0, 0.75,
    size=14, bold=True, color=C_GOLD1, align=PP_ALIGN.CENTER)
rect(sl, 0, H-0.35, W, 0.35, C_GOLD2)

# ═══════════════════════════════════════════════════════════════
# 저장
# ═══════════════════════════════════════════════════════════════
out = "/Users/hskim/Documents/repo/iitp/IITP_2026_연구사업계획_v2.pptx"
prs.save(out)
print(f"저장 완료 → {out}")
